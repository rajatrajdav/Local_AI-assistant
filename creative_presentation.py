"""
Creative Presentation Module (v3.0) — Multi-Theme Engine
=========================================================
Three professionally designed presentation themes with intelligent
auto-selection based on topic keywords.

THEMES:
  1. "corporate_edge"  — Dark charcoal, teal & lime green, geometric
  2. "wanderlust"      — Pastel blue/cream, rounded frames, airy
  3. "artistic_pitch"  — Vibrant fluid colors, black banners, bold

Each theme produces a distinctly styled presentation matching the
reference designs described by the user.
"""

import os
import time
import subprocess
import random
from datetime import datetime
from typing import List, Dict, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap
from lxml import etree

from pexels_integration import (
    get_pexels_client,
    download_background_image,
    get_color_palette_from_image,
    create_background_with_overlay,
    MEDIA_DIR,
)

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "generated_files")
os.makedirs(OUTPUT_DIR, exist_ok=True)


# ══════════════════════════════════════════════════════════════════════════════
# THEME DEFINITIONS
# ══════════════════════════════════════════════════════════════════════════════

class Theme:
    """A complete presentation theme with colors, fonts, shapes and layouts."""

    def __init__(self, name, display_name, description, colors, fonts, layout, shapes, slide_bg_color, title_slide_style):
        self.name = name
        self.display_name = display_name
        self.description = description
        self.colors = colors          # dict of color roles
        self.fonts = fonts            # dict: heading, body, accent
        self.layout = layout          # dict: margins, spacing, etc.
        self.shapes = shapes          # shape preferences
        self.slide_bg_color = slide_bg_color
        self.title_slide_style = title_slide_style


# Keywords that help auto-select the best theme for a topic
THEME_KEYWORDS = {
    "corporate_edge": [
        "business", "corporate", "finance", "technology", "software", "startup",
        "enterprise", "strategy", "consulting", "management", "analytics",
        "cybersecurity", "blockchain", "data science", "investor", "pitch deck",
        "saas", "b2b", "professional", "executive", "industry", "market",
        "economic", "corporate governance", "risk management", "audit",
        "compliance", "supply chain", "operations", "sales", "marketing",
    ],
    "wanderlust": [
        "travel", "tourism", "nature", "landscape", "culture", "photography",
        "lifestyle", "wellness", "yoga", "food", "cooking", "fashion",
        "design", "art", "creative", "furniture", "interior", "architecture",
        "garden", "outdoor", "adventure", "exploration", "vacation",
        "holiday", "beach", "mountain", "forest", "wildlife", "sustainability",
        "environment", "education", "learning", "school", "university",
    ],
    "artistic_pitch": [
        "fashion show", "art gallery", "exhibition", "creative agency",
        "portfolio", "branding", "advertising", "media", "entertainment",
        "music", "film", "theater", "performance", "digital art", "animation",
        "graphic design", "ui ux", "product design", "innovation lab",
        "creative director", "artistic", "modern art", "contemporary",
        "abstract", "visual arts", "multimedia", "photography portfolio",
    ],
}


def auto_select_theme(topic: str) -> str:
    """
    Automatically pick the best theme based on the presentation topic.
    Uses keyword matching with scoring.
    """
    topic_lower = topic.lower()
    scores = {"corporate_edge": 0, "wanderlust": 0, "artistic_pitch": 0}

    for theme, keywords in THEME_KEYWORDS.items():
        for kw in keywords:
            if kw in topic_lower:
                scores[theme] += 2
            # Also check word-by-word
            for word in topic_lower.split():
                if word == kw or word.startswith(kw) or kw.startswith(word):
                    scores[theme] += 1

    # If no match, use heuristics
    if max(scores.values()) == 0:
        # Default: corporate for tech/finance words, wanderlust for lifestyle, artistic for creative
        if any(w in topic_lower for w in ["technology", "software", "data", "business", "finance"]):
            return "corporate_edge"
        elif any(w in topic_lower for w in ["travel", "nature", "food", "design", "art"]):
            return "wanderlust"
        else:
            return "corporate_edge"  # safe default

    best_theme = max(scores, key=scores.get)
    return best_theme


# ══════════════════════════════════════════════════════════════════════════════
# THEME 1: CORPORATE EDGE
# ══════════════════════════════════════════════════════════════════════════════
# Dark charcoal gray bg, sharp teal & lime green accents, geometric frames,
# hourglass photo shapes, clean sans-serif, numbered lists, data blocks, icons.

CORPORATE_EDGE = Theme(
    name="corporate_edge",
    display_name="Corporate Edge",
    description="Dark charcoal with teal & lime green — professional, modern, high-impact",
    colors={
        "bg_dark":       (30, 32, 35),     # Charcoal gray
        "bg_medium":     (40, 42, 46),     # Slightly lighter
        "accent1":       (0, 173, 181),    # Teal
        "accent2":       (134, 224, 0),    # Lime green
        "accent3":       (255, 255, 255),  # White
        "text_primary":  (255, 255, 255),  # White text
        "text_secondary": (180, 190, 200), # Light gray
        "text_accent":   (0, 173, 181),    # Teal text
        "box_bg":        (38, 40, 44),     # Dark card bg
        "box_border":    (0, 173, 181),    # Teal border
        "number_color":  (134, 224, 0),    # Lime green numbers
    },
    fonts={
        "heading": "Calibri",
        "body":    "Calibri",
        "accent":  "Calibri Light",
        "mono":    "Consolas",
    },
    layout={
        "title_size":   44,
        "subtitle_size": 22,
        "heading_size":  30,
        "body_size":     16,
        "bullet_size":   14,
        "number_size":   48,
        "margin_left":   0.6,
        "margin_right":  0.6,
        "margin_top":    0.4,
        "content_left":  0.7,
        "content_width": 7.2,
        "spacing":       10,
    },
    shapes={
        "title_shape": MSO_SHAPE.RECTANGLE,
        "content_shape": MSO_SHAPE.ROUNDED_RECTANGLE,
        "accent_line": True,
        "use_geometric_frames": True,
    },
    slide_bg_color=(30, 32, 35),
    title_slide_style="bold_center",
)

# ══════════════════════════════════════════════════════════════════════════════
# THEME 2: WANDERLUST
# ══════════════════════════════════════════════════════════════════════════════
# Soft pastel blue & cream gradient bg, white circular line patterns,
# rounded rectangle image frames, airy bright photography, bold modern typography,
# spacious layout, organized grid system.

WANDERLUST = Theme(
    name="wanderlust",
    display_name="Wanderlust",
    description="Pastel blue & cream — airy, bright, spacious, friendly",
    colors={
        "bg_dark":       (200, 215, 230),  # Soft blue-gray (for gradient start)
        "bg_medium":     (240, 235, 220),  # Cream (for gradient end)
        "accent1":       (100, 150, 200),  # Soft blue
        "accent2":       (200, 180, 150),  # Warm beige
        "accent3":       (255, 255, 255),  # White
        "text_primary":  (40, 45, 55),     # Dark navy text
        "text_secondary": (100, 110, 120), # Medium gray
        "text_accent":   (80, 140, 190),   # Soft blue accent
        "box_bg":        (255, 255, 255),  # White cards
        "box_border":    (200, 215, 230),  # Soft blue border
        "number_color":  (100, 150, 200),  # Soft blue numbers
    },
    fonts={
        "heading": "Calibri Light",
        "body":    "Calibri",
        "accent":  "Calibri",
        "mono":    "Consolas",
    },
    layout={
        "title_size":   40,
        "subtitle_size": 20,
        "heading_size":  28,
        "body_size":     16,
        "bullet_size":   14,
        "number_size":   42,
        "margin_left":   0.7,
        "margin_right":  0.7,
        "margin_top":    0.5,
        "content_left":  0.5,
        "content_width": 7.5,
        "spacing":       12,
    },
    shapes={
        "title_shape": MSO_SHAPE.ROUNDED_RECTANGLE,
        "content_shape": MSO_SHAPE.ROUNDED_RECTANGLE,
        "accent_line": True,
        "use_circular_patterns": True,
        "rounded_corners": True,
    },
    slide_bg_color=(240, 235, 220),
    title_slide_style="centered_elegant",
)

# ══════════════════════════════════════════════════════════════════════════════
# THEME 3: ARTISTIC PITCH
# ══════════════════════════════════════════════════════════════════════════════
# Vibrant high-contrast abstract fluid colors (deep pink, magenta, royal blue),
# black sections, crisp white banners, large minimalist numbers, bold layout.

ARTISTIC_PITCH = Theme(
    name="artistic_pitch",
    display_name="Artistic Pitch",
    description="Vibrant fluid colors, black banners, bold & artistic",
    colors={
        "bg_dark":       (10, 5, 15),      # Near black with purple tint
        "bg_medium":     (25, 10, 35),     # Dark purple
        "accent1":       (220, 20, 100),   # Deep pink
        "accent2":       (180, 30, 140),   # Magenta
        "accent3":       (50, 60, 180),    # Royal blue
        "text_primary":  (255, 255, 255),  # White text
        "text_secondary": (200, 180, 200), # Light purple-gray
        "text_accent":   (255, 80, 140),   # Pink accent text
        "box_bg":        (15, 10, 25),     # Very dark purple
        "box_border":    (220, 20, 100),   # Pink border
        "number_color":  (255, 255, 255),  # White numbers
    },
    fonts={
        "heading": "Calibri",
        "body":    "Calibri",
        "accent":  "Calibri Light",
        "mono":    "Consolas",
    },
    layout={
        "title_size":   50,
        "subtitle_size": 24,
        "heading_size":  34,
        "body_size":     17,
        "bullet_size":   15,
        "number_size":   60,
        "margin_left":   0.5,
        "margin_right":  0.5,
        "margin_top":    0.3,
        "content_left":  0.6,
        "content_width": 6.8,
        "spacing":       8,
    },
    shapes={
        "title_shape": MSO_SHAPE.RECTANGLE,
        "content_shape": MSO_SHAPE.RECTANGLE,
        "accent_line": True,
        "use_color_blocks": True,
        "bold_banners": True,
    },
    slide_bg_color=(10, 5, 15),
    title_slide_style="dramatic_center",
)

# Registry of all themes
THEMES = {
    "corporate_edge": CORPORATE_EDGE,
    "wanderlust": WANDERLUST,
    "artistic_pitch": ARTISTIC_PITCH,
}


def get_theme(theme_name: str) -> Theme:
    """Get a theme by name, with fallback to corporate_edge."""
    return THEMES.get(theme_name, CORPORATE_EDGE)


# ══════════════════════════════════════════════════════════════════════════════
# OOXML HELPERS — for fills, borders, and effects python-pptx can't do natively
# ══════════════════════════════════════════════════════════════════════════════

def _apply_solid_fill(shape, r: int, g: int, b: int):
    """Apply a solid (opaque) fill to a shape via OOXML."""
    sp = shape.shape._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        spPr = etree.SubElement(sp, qn("p:spPr"))
    for tag in ["a:noFill", "a:solidFill", "a:gradFill", "a:pattFill", "a:blipFill"]:
        for child in spPr.findall(qn(tag)):
            spPr.remove(child)
    hex_color = f"{r:02X}{g:02X}{b:02X}"
    solid_fill = etree.SubElement(spPr, qn("a:solidFill"))
    srgb_clr = etree.SubElement(solid_fill, qn("a:srgbClr"))
    srgb_clr.set("val", hex_color)


def _apply_soft_gradient_fill(shape, color1: Tuple[int,int,int], color2: Tuple[int,int,int]):
    """Apply a linear gradient fill (top-to-bottom)."""
    sp = shape.shape._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        spPr = etree.SubElement(sp, qn("p:spPr"))
    for tag in ["a:noFill", "a:solidFill", "a:gradFill", "a:pattFill", "a:blipFill"]:
        for child in spPr.findall(qn(tag)):
            spPr.remove(child)

    grad = etree.SubElement(spPr, qn("a:gradFill"))
    gs_lst = etree.SubElement(grad, qn("a:gsLst"))

    # Color at 0%
    gs1 = etree.SubElement(gs_lst, qn("a:gs"))
    gs1.set("pos", "0")
    c1 = etree.SubElement(gs1, qn("a:srgbClr"))
    c1.set("val", f"{color1[0]:02X}{color1[1]:02X}{color1[2]:02X}")

    # Color at 100%
    gs2 = etree.SubElement(gs_lst, qn("a:gs"))
    gs2.set("pos", "100000")
    c2 = etree.SubElement(gs2, qn("a:srgbClr"))
    c2.set("val", f"{color2[0]:02X}{color2[1]:02X}{color2[2]:02X}")

    # Linear path
    lin = etree.SubElement(grad, qn("a:lin"))
    lin.set("ang", "5400000")   # 90 degrees (top to bottom)
    lin.set("scaled", "1")


def _apply_blended_fill(shape, r: int, g: int, b: int, alpha_percent: int):
    """Apply semi-transparent fill (glassmorphism effect)."""
    sp = shape.shape._element
    hex_color = f"{r:02X}{g:02X}{b:02X}"
    alpha_val = int((alpha_percent / 100) * 100000)
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        spPr = etree.SubElement(sp, qn("p:spPr"))
    for tag in ["a:noFill", "a:solidFill", "a:gradFill", "a:pattFill", "a:blipFill"]:
        for child in spPr.findall(qn(tag)):
            spPr.remove(child)
    solid_fill = etree.SubElement(spPr, qn("a:solidFill"))
    srgb_clr = etree.SubElement(solid_fill, qn("a:srgbClr"))
    srgb_clr.set("val", hex_color)
    alpha_elem = etree.SubElement(srgb_clr, qn("a:alpha"))
    alpha_elem.set("val", str(100000 - alpha_val))


def _remove_shape_border(shape):
    """Remove border from a shape."""
    try:
        shape.line.fill.background()
    except Exception:
        pass


def _set_shape_border(shape, r: int, g: int, b: int, width: float = 1.5):
    """Set a colored border on a shape via OOXML."""
    sp = shape.shape._element
    hex_color = f"{r:02X}{g:02X}{b:02X}"
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        spPr = etree.SubElement(sp, qn("p:spPr"))
    # Remove existing ln
    for child in spPr.findall(qn("a:ln")):
        spPr.remove(child)
    # Add new ln
    ln = etree.SubElement(spPr, qn("a:ln"))
    ln.set("w", str(int(width * 12700)))  # EMU
    sf = etree.SubElement(ln, qn("a:solidFill"))
    sc = etree.SubElement(sf, qn("a:srgbClr"))
    sc.set("val", hex_color)


def _set_shape_shadow(shape, r: int, g: int, b: int, blur_radius: int = 100000):
    """Add a soft shadow to a shape."""
    sp = shape.shape._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        spPr = etree.SubElement(sp, qn("p:spPr"))
    # Remove existing shadow
    for child in spPr.findall(qn("a:effectLst")):
        spPr.remove(child)
    effect_lst = etree.SubElement(spPr, qn("a:effectLst"))
    outer_shdw = etree.SubElement(effect_lst, qn("a:outerShdw"))
    outer_shdw.set("blurRad", str(blur_radius))
    outer_shdw.set("dist", "50000")
    outer_shdw.set("dir", "2700000")
    outer_shdw.set("algn", "tl")
    outer_shdw.set("rotWithShape", "0")
    srgb = etree.SubElement(outer_shdw, qn("a:srgbClr"))
    srgb.set("val", f"{r:02X}{g:02X}{b:02X}")
    alpha = etree.SubElement(srgb, qn("a:alpha"))
    alpha.set("val", "50000")  # ~50%


def _add_circular_pattern(slide, color: Tuple[int,int,int], count: int = 12):
    """Add subtle decorative circles to a slide (for Wanderlust theme)."""
    slide_width = Inches(13.333)
    slide_height = Inches(7.5)
    for i in range(count):
        r = random.randint(60, 200)
        x = random.randint(0, int(13.333 * 914400 - r))
        y = random.randint(0, int(7.5 * 914400 - r))
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Emu(x), Emu(y), Emu(r), Emu(r),
        )
        _apply_blended_fill(circle, color[0], color[1], color[2], 88)  # ~12% opacity
        _remove_shape_border(circle)
        # Send to back by moving XML order
        sp = circle.shape._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)  # Behind content


def _add_geometric_frame(slide, left, top, width, height, color: Tuple[int,int,int]):
    """Add a geometric frame shape (for Corporate Edge theme)."""
    # Hourglass-like geometric shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.PENTAGON, left, top, width, height,
    )
    _apply_blended_fill(shape, color[0], color[1], color[2], 80)  # Very subtle
    _set_shape_border(shape, color[0], color[1], color[2], 0.5)
    return shape


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE DESIGNER — Multi-Theme
# ══════════════════════════════════════════════════════════════════════════════

class CreativeSlideDesigner:
    """
    Designs slides using the specified theme.
    Supports Corporate Edge, Wanderlust, and Artistic Pitch.
    """

    def __init__(self, theme: Theme):
        self.theme = theme
        self.pexels_client = get_pexels_client()
        self.slide_width = Inches(13.333)
        self.slide_height = Inches(7.5)

    def get_text_color(self, bg_colors: Dict) -> Tuple[int, int, int]:
        if bg_colors.get("is_dark_background", True):
            return self.theme.colors["text_primary"]
        return self.theme.colors["text_primary"]

    def rgb_to_pptx_color(self, rgb: Tuple[int, int, int]):
        return RGBColor(rgb[0], rgb[1], rgb[2])

    def add_slide_background(self, slide, base_color: Tuple[int,int,int] = None):
        """Fill the slide with the theme's background color/gradient."""
        if base_color is None:
            base_color = self.theme.slide_bg_color

        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            width=self.slide_width, height=self.slide_height,
        )
        bg.name = "SlideBackground"

        # Apply gradient or solid based on theme
        if self.theme.name == "wanderlust":
            # Soft pastel gradient
            c1 = (self.theme.colors["bg_dark"][0], self.theme.colors["bg_dark"][1], self.theme.colors["bg_dark"][2])
            c2 = (self.theme.colors["bg_medium"][0], self.theme.colors["bg_medium"][1], self.theme.colors["bg_medium"][2])
            _apply_soft_gradient_fill(bg, c1, c2)
            # Add decorative circles
            _add_circular_pattern(slide, self.theme.colors["accent1"], count=8)
        elif self.theme.name == "artistic_pitch":
            # Dark dramatic background
            _apply_solid_fill(bg, base_color[0], base_color[1], base_color[2])
        else:
            # Corporate Edge — solid charcoal
            _apply_solid_fill(bg, base_color[0], base_color[1], base_color[2])

        _remove_shape_border(bg)
        # Send to back
        sp = bg.shape._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)

    def add_accent_bar(self, slide, left, top, width, height):
        """Add a colored accent bar/shape."""
        color = self.theme.colors["accent1"]
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height,
        )
        _apply_solid_fill(bar, color[0], color[1], color[2])
        _remove_shape_border(bar)
        return bar

    def add_accent_line(self, slide, left, top, width):
        """Add a thin accent line."""
        color = self.theme.colors["accent2"]
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, Pt(3),
        )
        _apply_solid_fill(line, color[0], color[1], color[2])
        _remove_shape_border(line)

    def add_number_badge(self, slide, number: int, left, top):
        """Add a large numbered badge (for numbered lists)."""
        c = self.theme.colors["number_color"]
        num = slide.shapes.add_textbox(left, top, Inches(0.8), Inches(0.8))
        tf = num.text_frame
        p = tf.paragraphs[0]
        p.text = str(number)
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(self.theme.layout["number_size"])
        p.font.bold = True
        p.font.color.rgb = self.rgb_to_pptx_color(c)
        p.font.name = self.theme.fonts["heading"]

    def add_title_slide(self, slide, title: str, subtitle: str, bg_image: str = None, palette: Dict = None):
        """Create the title slide in the theme's style."""
        self.add_slide_background(slide)

        # Add background image if available (as overlay)
        if bg_image and os.path.exists(bg_image):
            try:
                slide.shapes.add_picture(
                    bg_image, 0, 0,
                    width=self.slide_width, height=self.slide_height,
                )
                # Re-add background on top with low opacity for readability
                overlay = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 0, 0,
                    width=self.slide_width, height=self.slide_height,
                )
                bg_c = self.theme.slide_bg_color
                _apply_blended_fill(overlay, bg_c[0], bg_c[1], bg_c[2], 35)
                _remove_shape_border(overlay)
            except Exception:
                pass

        if self.theme.name == "corporate_edge":
            # Bold centered title with teal accent bar
            self.add_accent_bar(slide, Inches(4.5), Inches(2.8), Inches(4.3), Pt(4))

            tbox = slide.shapes.add_textbox(
                Inches(1.5), Inches(3.0), Inches(10.3), Inches(1.8),
            )
            tf = tbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(self.theme.layout["title_size"])
            p.font.bold = True
            p.font.color.rgb = self.rgb_to_pptx_color(self.theme.colors["text_primary"])
            p.font.name = self.theme.fonts["heading"]

            if subtitle:
                sbox = slide.shapes.add_textbox(
                    Inches(2.5), Inches(4.8), Inches(8.3), Inches(1.2),
                )
                stf = sbox.text_frame
                stf.word_wrap = True
                sp = stf.paragraphs[0]
                sp.text = subtitle
                sp.alignment = PP_ALIGN.CENTER
                sp.font.size = Pt(self.theme.layout["subtitle_size"])
                sp.font.color.rgb = self.rgb_to_pptx_color(self.theme.colors["text_secondary"])
                sp.font.name = self.theme.fonts["body"]

        elif self.theme.name == "wanderlust":
            # Centered elegant with soft blue header
            tbox = slide.shapes.add_textbox(
                Inches(2.0), Inches(2.5), Inches(9.3), Inches(1.5),
            )
            tf = tbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(self.theme.layout["title_size"])
            p.font.bold = False
            p.font.color.rgb = self.rgb_to_pptx_color(self.theme.colors["text_accent"])
            p.font.name = self.theme.fonts["heading"]

            # Decorative line below title
            self.add_accent_line(slide, Inches(5.5), Inches(4.0), Inches(2.3))

            if subtitle:
                sbox = slide.shapes.add_textbox(
                    Inches(2.5), Inches(4.3), Inches(8.3), Inches(1.0),
                )
                stf = sbox.text_frame
                stf.word_wrap = True
                sp = stf.paragraphs[0]
                sp.text = subtitle
                sp.alignment = PP_ALIGN.CENTER
                sp.font.size = Pt(self.theme.layout["subtitle_size"])
                sp.font.color.rgb = self.rgb_to_pptx_color(self.theme.colors["text_secondary"])
                sp.font.name = self.theme.fonts["body"]

        else:  # artistic_pitch
            # Dramatic centered — large bold white on dark
            self.add_accent_bar(slide, 0, Inches(3.0), self.slide_width, Pt(2))

            tbox = slide.shapes.add_textbox(hii 
                                            i
                Inches(1.0), Inches(1.8), Inches(11.3), Inches(2.2),
            )
            tf = tbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(self.theme.layout["title_size"])
            p.font.bold = True
            p.font.color.rgb = self.rgb_to_pptx_color(self.theme.colors["text_primary"])
            p.font.name = self.theme.fonts["heading"]

            if subtitle:
                sbox = slide.shapes.add_textbox(
                    Inches(2.0), Inches(4.2), Inches(9.3), Inches(1.0),
                )
                stf = sbox.text_frame
                stf.word_wrap = True
                sp = stf.paragraphs[0]
                sp.text = subtitle
                sp.alignment = PP_ALIGN.CENTER
                sp.font.size = Pt(self.theme.layout["subtitle_size"])
                sp.font.color.rgb = self.rgb_to_pptx_color(self.theme.colors["text_accent"])
                sp.font.name = self.theme.fonts["accent"]

    def add_content_box(
        self,
        slide,
        title: str,
        content: str,
        position: str = "right",
        slide_index: int = 0,
    ):
        """
        Add a styled content box matching the active theme.
        """
        c = self.theme.colors
        layout = self.theme.layout

        positions = {
            "left":   (Inches(0.4),  Inches(1.0), Inches(5.8),  Inches(6.0)),
            "right":  (Inches(7.0),  Inches(1.0), Inches(5.8),  Inches(6.0)),
            "center": (Inches(1.5),  Inches(1.0), Inches(10.3), Inches(6.0)),
        }
        left, top, width, height = positions.get(position, positions["right"])

        # ── Background card ──────────────────────────────────────────────────
        if self.theme.name == "wanderlust":
            # White rounded card with soft shadow
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
            )
            _apply_solid_fill(card, c["box_bg"][0], c["box_bg"][1], c["box_bg"][2])
            _set_shape_border(card, c["box_border"][0], c["box_border"][1], c["box_border"][2], 0.8)
            _set_shape_shadow(card, 150, 150, 180, blur_radius=80000)
            # Adjust corner rounding
            spPr = card.shape._element.find(qn("p:spPr"))
            if spPr is not None:
                prstGeom = spPr.find(qn("a:prstGeom"))
                if prstGeom is not None:
                    prstGeom.set("prst", "roundRect")

        elif self.theme.name == "artistic_pitch":
            # Dark panel with pink border
            card = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, width, height,
            )
            _apply_solid_fill(card, c["box_bg"][0], c["box_bg"][1], c["box_bg"][2])
            _set_shape_border(card, c["box_border"][0], c["box_border"][1], c["box_border"][2], 1.0)

        else:  # corporate_edge
            # Dark card with teal accent bar on left edge
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
            )
            _apply_solid_fill(card, c["box_bg"][0], c["box_bg"][1], c["box_bg"][2])
            _set_shape_border(card, c["box_border"][0], c["box_border"][1], c["box_border"][2], 0.5)

            # Small accent bar on left edge of card
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, top, Pt(4), height,
            )
            _apply_solid_fill(bar, c["accent1"][0], c["accent1"][1], c["accent1"][2])
            _remove_shape_border(bar)

        # ── Text content ─────────────────────────────────────────────────────
        tf = card.text_frame if self.theme.name == "wanderlust" else (
            slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.15), width - Inches(0.6), height - Inches(0.3))
        ).text_frame

        if self.theme.name != "wanderlust":
            tf = card.text_frame

        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.2)
        tf.margin_bottom = Inches(0.15)

        # Title paragraph
        title_p = tf.paragraphs[0]
        title_p.text = title
        title_p.alignment = PP_ALIGN.LEFT
        title_p.font.size = Pt(layout["heading_size"])
        title_p.font.bold = True
        title_p.font.name = self.theme.fonts["heading"]
        title_p.font.color.rgb = self.rgb_to_pptx_color(c["text_accent"])

        # Separator line
        sep = tf.add_paragraph()
        sep.text = "─" * 30
        sep.font.size = Pt(6)
        sep.font.color.rgb = self.rgb_to_pptx_color(c["text_secondary"])
        sep.space_after = Pt(6)

        # Content bullets
        lines = content.split("\n") if content else []
        for i, line in enumerate(lines):
            stripped = line.strip()
            if not stripped:
                continue
            if stripped.startswith(("•", "-", "*", "→", "▸")):
                stripped = stripped[1:].strip()

            # Choose bullet character based on theme
            if self.theme.name == "corporate_edge":
                bullet = "▸"
            elif self.theme.name == "wanderlust":
                bullet = "◦"
            else:  # artistic_pitch
                bullet = "—"

            bp = tf.add_paragraph()
            bp.text = f"{bullet}  {stripped}"
            bp.alignment = PP_ALIGN.LEFT
            bp.font.size = Pt(layout["bullet_size"])
            bp.font.name = self.theme.fonts["body"]
            bp.font.color.rgb = self.rgb_to_pptx_color(c["text_primary"])
            bp.space_after = Pt(layout["spacing"])
            bp.space_before = Pt(3)


# ══════════════════════════════════════════════════════════════════════════════
# PRESENTATION CREATOR — Multi-Theme
# ══════════════════════════════════════════════════════════════════════════════

class CreativePresentationCreator:
    """Creates presentations using the specified or auto-detected theme."""

    def __init__(self, theme_name: str = None):
        """
        Initialize with a theme. If None, will be auto-selected from the topic.
        """
        self.theme_name = theme_name
        self.theme = None
        self.designer = None
        self.pexels_client = get_pexels_client()

    def create_presentation(
        self,
        title: str,
        slides_content: List[Dict],
        filename: str = None,
        show_progress: bool = True,
        theme: str = None,
    ) -> str:
        """
        Create a presentation with the best-matched theme.

        Args:
            title: Presentation title
            slides_content: List of slide dicts with title, content, image_query
            filename: Output filename
            show_progress: Print progress to console
            theme: Force a specific theme, or None for auto-detect
        """
        # ── Determine theme ─────────────────────────────────────────────────
        theme_name = theme or self.theme_name
        if theme_name and theme_name in THEMES:
            self.theme = get_theme(theme_name)
        else:
            # Auto-detect from title
            detected = auto_select_theme(title)
            self.theme = get_theme(detected)
            theme_name = detected

        self.designer = CreativeSlideDesigner(self.theme)

        if filename is None:
            ts = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"{theme_name}_presentation_{ts}.pptx"
        if not filename.endswith(".pptx"):
            filename += ".pptx"
        filepath = os.path.join(OUTPUT_DIR, filename)

        if show_progress:
            print(f"\n  {'=' * 56}")
            print(f"  🎨 THEME: {self.theme.display_name}")
            print(f"  📝 {self.theme.description}")
            print(f"  {'=' * 56}")
            print(f"\n  📊 Creating: {title}")
            print(f"  🖼  Slides: {len(slides_content) + 1}")
            time.sleep(1.5)

        prs = Presentation()
        prs.slide_width = self.designer.slide_width
        prs.slide_height = self.designer.slide_height

        # ── TITLE SLIDE ─────────────────────────────────────────────────────
        if show_progress:
            print(f"\n  📄 Creating Title Slide...")
            time.sleep(0.5)

        title_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        # Try to get a background image for the title slide
        if show_progress:
            print(f"     📷 Searching for relevant background image...")
            time.sleep(0.3)

        bg_image = download_background_image(f"{title} background", 0)
        palette = None
        if bg_image and os.path.exists(bg_image):
            palette = get_color_palette_from_image(bg_image)
            if show_progress:
                print(f"     ✅ Background found for title slide")

        subtitle = (
            f"A Comprehensive Professional Overview\n"
            f"{datetime.now().strftime('%B %Y')}"
        )
        self.designer.add_title_slide(title_slide, title, subtitle, bg_image, palette)

        if show_progress:
            print(f"     ✓ Title slide created")
            time.sleep(1.0)

        # ── CONTENT SLIDES ─────────────────────────────────────────────────
        if show_progress:
            print(f"\n  ─── Creating {len(slides_content)} content slides ───")
            time.sleep(0.8)

        for i, slide_data in enumerate(slides_content, 1):
            slide_title   = slide_data.get("title", f"Slide {i}")
            slide_content = slide_data.get("content", "")
            image_query   = slide_data.get("image_query", slide_title)
            text_position = slide_data.get("text_position", "right")

            if show_progress:
                print(f"\n  ════════════════════════════════════════════")
                print(f"  🎨 Slide {i}/{len(slides_content)}: {slide_title}")
                print(f"  ════════════════════════════════════════════")
                time.sleep(0.5)

            # Step 1: Get background
            if show_progress:
                print(f"     ⏳ Searching Pexels for '{image_query}'...")
                time.sleep(0.3)

            bg_image = download_background_image(image_query, i)
            palette = None
            if bg_image and os.path.exists(bg_image):
                palette = get_color_palette_from_image(bg_image)

            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

            # Step 2: Apply themed background
            if show_progress:
                print(f"     🎨 Applying {self.theme.display_name} theme styling...")
                time.sleep(0.3)

            self.designer.add_slide_background(slide)

            # Step 3: Add background image as decorative overlay if available
            if bg_image and os.path.exists(bg_image):
                try:
                    # Add image as full-bleed with very subtle opacity via overlay
                    slide.shapes.add_picture(
                        bg_image, 0, 0,
                        width=self.designer.slide_width,
                        height=self.designer.slide_height,
                    )
                    # Cover with very transparent overlay for readability
                    img_overlay = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, 0, 0,
                        width=self.designer.slide_width,
                        height=self.designer.slide_height,
                    )
                    bg_c = self.theme.slide_bg_color
                    _apply_blended_fill(img_overlay, bg_c[0], bg_c[1], bg_c[2], 50)
                    _remove_shape_border(img_overlay)
                    if show_progress:
                        print(f"     🖼  Background image blended into slide")
                except Exception as e:
                    if show_progress:
                        print(f"     ⚠ Image overlay skipped: {e}")

            # Step 4: Add title bar
            title_bar_top = Inches(0.3)
            title_bar = slide.shapes.add_textbox(
                Inches(0.6), title_bar_top, Inches(12.0), Inches(0.7),
            )
            tf = title_bar.text_frame
            tf.word_wrap = True
            tp = tf.paragraphs[0]
            tp.text = f"  {slide_title}"
            tp.font.size = Pt(self.theme.layout["heading_size"])
            tp.font.bold = True
            tp.font.color.rgb = self.designer.rgb_to_pptx_color(self.theme.colors["text_accent"])
            tp.font.name = self.theme.fonts["heading"]

            # Step 5: Add content
            if show_progress:
                print(f"     ✍️ Adding styled content panel...")
                time.sleep(0.2)

            self.designer.add_content_box(
                slide, slide_title, slide_content,
                position=text_position, slide_index=i,
            )

            if show_progress:
                print(f"     ✅ Slide {i} complete!")
                time.sleep(0.6)

            if show_progress and i < len(slides_content):
                print(f"     ⏳ Preparing next slide...")
                time.sleep(0.4)

        # ── SAVE ──────────────────────────────────────────────────────────────
        if show_progress:
            print(f"\n  ════════════════════════════════════════════")
            print(f"  💾 Saving presentation to disk...")
            time.sleep(0.3)

        prs.save(filepath)

        if show_progress:
            print(f"     ✓ File saved: {filepath}")
            print(f"\n  ✅ Presentation complete! {len(slides_content) + 1} slides.")
            print(f"  🎨 Theme: {self.theme.display_name}")
            print(f"  📁 Filename: {filename}")
            time.sleep(0.3)

        # Auto-open
        _auto_open_file(filepath)
        return filepath


def _auto_open_file(filepath: str):
    """Open the generated file with the OS default application."""
    try:
        if os.name == "nt":
            os.startfile(filepath)
        elif hasattr(os, "uname") and os.uname().sysname == "Darwin":
            subprocess.Popen(["open", filepath])
        else:
            subprocess.Popen(["xdg-open", filepath])
        print(f"  🚀 Presentation opened: {os.path.basename(filepath)}")
    except Exception as e:
        print(f"  [Auto-open failed]: {e}")
        print(f"  📁 Open manually: {filepath}")


# ══════════════════════════════════════════════════════════════════════════════
# PUBLIC CONVENIENCE FUNCTION (backward-compatible signature)
# ══════════════════════════════════════════════════════════════════════════════

def create_creative_presentation(
    title: str,
    slides_content: List[Dict],
    filename: str = None,
    show_progress: bool = True,
    theme: str = None,
) -> str:
    """
    Create a presentation with automatic theme selection.

    Args:
        title: Presentation title (used for theme auto-detection)
        slides_content: List of slide dicts
        filename: Optional output filename
        show_progress: Show progress in console
        theme: Force a theme ("corporate_edge", "wanderlust", "artistic_pitch"),
               or None for auto-detect

    Returns:
        Path to the saved .pptx file
    """
    creator = CreativePresentationCreator(theme_name=theme)
    return creator.create_presentation(title, slides_content, filename, show_progress, theme=theme)


# ══════════════════════════════════════════════════════════════════════════════
# TEST RUNNER
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    print("Testing Creative Presentation System v3.0 — Multi-Theme Engine")
    print("=" * 60)

    # Test slides
    test_slides = [
        {
            "title": "Introduction & Background",
            "content": (
                "Artificial Intelligence represents the simulation of human intelligence "
                "processes by computer systems, enabling machines to learn from data.\n"
                "The field encompasses machine learning, deep learning, neural networks, "
                "and natural language processing as its primary technical disciplines.\n"
                "AI systems today power applications across healthcare, finance, "
                "transportation, and manufacturing with unprecedented efficiency.\n"
                "The global AI market is projected to reach $1.8 trillion by 2030, "
                "reflecting its transformative impact across every major industry."
            ),
            "image_query": "artificial intelligence technology",
        },
        {
            "title": "Core Machine Learning Concepts",
            "content": (
                "Supervised learning trains models on labeled datasets to make "
                "predictions on unseen data with measurable accuracy.\n"
                "Unsupervised learning discovers hidden patterns within unlabeled "
                "data through clustering and dimensionality reduction techniques.\n"
                "Reinforcement learning enables agents to learn optimal behaviors "
                "through trial and error with reward-based feedback mechanisms.\n"
                "Transfer learning allows pre-trained models to be fine-tuned for "
                "new tasks, dramatically reducing training time and data requirements."
            ),
            "image_query": "neural network data science",
        },
        {
            "title": "Industry Applications",
            "content": (
                "Healthcare AI systems can detect diseases from medical imaging "
                "with accuracy rates exceeding human radiologists in controlled studies.\n"
                "Financial institutions deploy AI for fraud detection, algorithmic "
                "trading, and personalized banking recommendations at massive scale.\n"
                "Autonomous vehicles rely on AI perception systems that process "
                "real-time sensor data to navigate complex road environments safely.\n"
                "E-commerce platforms leverage AI recommendation engines that "
                "drive 35% of consumer purchases through personalized suggestions."
            ),
            "image_query": "business technology innovation",
        },
    ]

    # Test all three themes
    test_cases = [
        ("Artificial Intelligence: A Comprehensive Overview", test_slides, None),
        ("Travel & Lifestyle: Exploring New Horizons", test_slides, "wanderlust"),
        ("Creative Agency Portfolio: Brand Innovation", test_slides, "artistic_pitch"),
    ]

    for title, slides, theme in test_cases:
        print(f"\n{'─' * 60}")
        print(f"TEST: {title}")
        if theme:
            print(f"THEME: {theme}")
        else:
            auto = auto_select_theme(title)
            print(f"AUTO-SELECTED THEME: {auto}")
        print(f"{'─' * 60}")

        result = create_creative_presentation(
            title, slides, show_progress=True, theme=theme,
        )
        print(f"\n✓ Saved to: {result}")