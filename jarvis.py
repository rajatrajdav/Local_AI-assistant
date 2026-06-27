import os
import asyncio
import json
import time
import re
import numpy as np
import sounddevice as sd
from piper import PiperVoice
from dotenv import load_dotenv

# ============================================================d
# Load environment variables from .env (API keys, secrets)
# NEVER hardcode API keys — they stay in .env which is gitignored
# ============================================================
load_dotenv()

# ============================================================
# Console-mode UI stubs (ui_jarvis.py removed)
# These no-op functions allow jarvis.py to run in pure console mode
# without requiring the GUI/UI module.
# ============================================================
def _ui_update(**kwargs):
    pass

def _ui_message(msg):
    pass

def _ui_notification(title, msg_type="info"):
    pass

# ============================================================
# FIX #1 — HINDI TEXT NORMALIZER UTILITY
# Cleans and normalizes Hindi text before passing to TTS engine
# to prevent mispronunciation and garbled audio output.
# ============================================================

def normalize_hindi_text(text: str) -> str:
    """
    Normalize Hindi (Devanagari) text for clean TTS pronunciation.
    Removes emojis, strips English mixed noise, normalizes punctuation.
    """
    # Remove emojis and unicode symbols that confuse TTS
    emoji_pattern = re.compile(
        "["
        "\U0001F600-\U0001F64F"
        "\U0001F300-\U0001F5FF"
        "\U0001F680-\U0001F6FF"
        "\U0001F1E0-\U0001F1FF"
        "\U00002500-\U00002BEF"
        "\U00002702-\U000027B0"
        "\U000024C2-\U0001F251"
        "]+",
        flags=re.UNICODE,
    )
    text = emoji_pattern.sub("", text)

    # Replace English punctuation with Hindi-friendly equivalents
    text = text.replace(".", "।")
    text = text.replace("!", "।")
    text = text.replace("?", "?")     # Keep ? — Piper handles it
    text = text.replace(",", ",")

    # Remove markdown symbols
    text = re.sub(r"[*_`#\[\]()]", "", text)

    # Collapse multiple spaces/newlines
    text = re.sub(r"\s+", " ", text).strip()

    # Remove any remaining non-Devanagari / non-punctuation characters
    # that may cause the TTS engine to stumble
    cleaned = ""
    for char in text:
        cp = ord(char)
        # Allow: Devanagari (0900–097F), spaces, common punctuation
        if (0x0900 <= cp <= 0x097F) or char in " ,।?-\n":
            cleaned += char
        elif char.isascii() and char.isprintable():
            # Keep ASCII letters/digits (mixed code is common in Indian speech)
            cleaned += char
    return cleaned.strip()


def normalize_english_text(text: str) -> str:
    """
    Light normalization for English TTS — removes markdown, emojis.
    """
    emoji_pattern = re.compile(
        "["
        "\U0001F600-\U0001F64F"
        "\U0001F300-\U0001F5FF"
        "\U0001F680-\U0001F6FF"
        "\U0001F1E0-\U0001F1FF"
        "\U00002500-\U00002BEF"
        "\U00002702-\U000027B0"
        "\U000024C2-\U0001F251"
        "]+",
        flags=re.UNICODE,
    )
    text = emoji_pattern.sub("", text)
    text = re.sub(r"[*_`#]", "", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text


# ============================================================
# Groq API integration for ultra-fast LLM responses
# ============================================================
try:
    from groq import Groq
    GROQ_API_KEY = os.getenv("GROQ_API_KEY")
    if not GROQ_API_KEY:
        print("⚠ GROQ_API_KEY not found in .env file")
        print("   Create a .env file with: GROQ_API_KEY=your_key_here")
        USE_GROQ = False
    else:
        groq_client = Groq(api_key=GROQ_API_KEY)
        USE_GROQ = True
        print("✓ Groq API connected - Ultra-fast mode enabled!")
except Exception as e:
    print(f"⚠ Groq API connection failed: {e}")
    USE_GROQ = False

# ============================================================
# JARVIS AGENT BRAIN — Google Gen AI SDK + Context Caching + Key Rotation
# Uses the new `google-genai` package (NOT deprecated `google.generativeai`)
# ============================================================
try:
    from google import genai
    from google.genai import types as genai_types
    USE_NEW_GENAI_SDK = True
except ImportError:
    USE_NEW_GENAI_SDK = False
    print("⚠ New google-genai SDK not installed. Run: pip install google-genai")

try:
    import groq as _groq_module
    GROQ_AVAILABLE = True
except ImportError:
    GROQ_AVAILABLE = False

class JarvisAgentBrain:
    """
    Heavy-duty coding agent core with:
    - Google API Key Rotation (seamlessly swaps keys on rate-limit)
    - Context Caching (upload codebase once, queries only consume new tokens)
    - Automatic Groq fallback when ALL Google keys are exhausted
    """

    def __init__(self):
        # ── Google API keys (up to 3 for rotation) ──
        self.google_keys = [
            os.getenv("GEMINI_API_KEY"),
            os.getenv("GEMINI_API_KEY_2"),
            os.getenv("GEMINI_API_KEY_3"),
        ]
        # Filter out None/empty keys
        self.google_keys = [k for k in self.google_keys if k and k.strip()]
        self.current_key_index = 0
        self.groq_key = os.getenv("GROQ_API_KEY")

        # ── Clients ──
        self.google_client = None
        self.groq_client = None
        self.model_name = os.getenv("GEMINI_MODEL", "gemini-1.5-flash")
        self.active_cache = None
        self.cache_ttl = int(os.getenv("CACHE_TTL_SECONDS", "7200"))

        # ── Initialize Google client ──
        if USE_NEW_GENAI_SDK and self.google_keys:
            try:
                self.google_client = genai.Client(api_key=self.google_keys[0])
                print(f"✓ Google Gen AI SDK connected — Key Rotation enabled ({len(self.google_keys)} key(s))")
            except Exception as e:
                print(f"⚠ Google Gen AI SDK init failed: {e}")
                self.google_client = None

        # ── Initialize Groq client ──
        if GROQ_AVAILABLE and self.groq_key:
            try:
                self.groq_client = _groq_module.Groq(api_key=self.groq_key)
            except Exception as e:
                print(f"⚠ Groq init failed: {e}")
                self.groq_client = None

        if self.google_client:
            print(f"  Model: {self.model_name} | Cache TTL: {self.cache_ttl}s")
        if self.groq_client:
            print("  Groq fallback: READY")

    # ────────────────────────────────────────────────────────
    # KEY ROTATION
    # ────────────────────────────────────────────────────────

    def _rotate_google_key(self):
        """Swap to next Google API key on rate-limit (429)."""
        if len(self.google_keys) <= 1:
            return False
        self.current_key_index = (self.current_key_index + 1) % len(self.google_keys)
        new_key = self.google_keys[self.current_key_index]
        print(f"[JarvisBrain] Rotating to Google API Key Index: {self.current_key_index}")
        try:
            self.google_client = genai.Client(api_key=new_key)
            return True
        except Exception:
            return False

    def google_available(self) -> bool:
        """Check if any Google key is currently usable."""
        return self.google_client is not None and len(self.google_keys) > 0

    def groq_available(self) -> bool:
        """Check if Groq client is usable."""
        return self.groq_client is not None

    # ────────────────────────────────────────────────────────
    # CONTEXT CACHING — Save massive token costs
    # ────────────────────────────────────────────────────────

    def update_codebase_cache(self, whole_codebase_text: str, ttl_seconds: int = None):
        """
        Upload and cache your entire codebase context.
        Subsequent queries will use the cached context, saving ~90% of input tokens.
        
        Requirements:
        - Minimum 2,048 tokens to create a cache
        - Default TTL: 2 hours (7200 seconds)
        - Cache is auto-extended on each query hit
        """
        if not self.google_client:
            print("[JarvisBrain] No Google client — cannot create cache.")
            return False

        ttl = ttl_seconds or self.cache_ttl
        
        # Estimate token count (rough: 4 chars ≈ 1 token)
        estimated_tokens = len(whole_codebase_text) // 4
        if estimated_tokens < 2048:
            print(f"[JarvisBrain] Codebase too small for caching ({estimated_tokens} est. tokens, need ≥2048). Skipping cache.")
            return False

        try:
            print(f"[JarvisBrain] Creating server-side context cache ({estimated_tokens} est. tokens, TTL={ttl}s)...")
            self.active_cache = self.google_client.caches.create(
                model=self.model_name,
                config=genai_types.CreateCachedContentConfig(
                    contents=[whole_codebase_text],
                    system_instruction=(
                        "You are Jarvis's elite coding core. "
                        "Use the provided codebase context to precisely handle "
                        "the user's software engineering requests."
                    ),
                    ttl=f"{ttl}s",
                ),
            )
            print(f"✓ Codebase cached! Cache ID: {self.active_cache.name}")
            return True
        except Exception as e:
            print(f"[JarvisBrain] Cache creation failed: {e}")
            return False

    # ────────────────────────────────────────────────────────
    # CORE QUERY ENGINE — Google → Rotation → Groq
    # ────────────────────────────────────────────────────────

    def _is_rate_limit_error(self, err_str: str) -> bool:
        """Detect rate-limit / quota errors to trigger key rotation."""
        err_lower = err_str.lower()
        patterns = [
            "429", "rate_limit", "rate limit", "quota", "insufficient_quota",
            "resource exhausted", "too many requests", "token expired",
            "unauthorized", "401", "403", "forbidden", "invalid api key",
        ]
        return any(p in err_lower for p in patterns)

    def ask_agent(self, prompt: str, system_instruction: str = None) -> str:
        """
        Send a query to the LLM with automatic:
        1. Context cache usage (if available)
        2. Google key rotation (on 429 errors)
        3. Groq emergency fallback
        
        Returns the response text.
        """
        # ── Try Google (with Context Cache) ──
        if self.google_available():
            config_kwargs = {}
            
            # Attach cached content to save tokens
            if self.active_cache:
                config_kwargs["cached_content"] = self.active_cache.name
            
            if system_instruction:
                config_kwargs["system_instruction"] = system_instruction

            # Try each Google key at most once
            for attempt in range(len(self.google_keys)):
                try:
                    response = self.google_client.models.generate_content(
                        model=self.model_name,
                        contents=prompt,
                        config=genai_types.GenerateContentConfig(**config_kwargs)
                    )

                    # Report cache hit
                    if response.usage_metadata and response.usage_metadata.cached_content_token_count:
                        saved = response.usage_metadata.cached_content_token_count
                        print(f"⚡ [Cache Hit!] Saved {saved} input tokens.")

                    return response.text

                except Exception as e:
                    err_str = str(e)
                    print(f"[Google attempt {attempt+1}] Error: {type(e).__name__}: {str(e)[:100]}")
                    
                    if self._is_rate_limit_error(err_str):
                        if self._rotate_google_key():
                            print(f"  → Rotated to key index {self.current_key_index}, retrying...")
                            continue
                        else:
                            print("  → No more keys to rotate to.")
                            break
                    else:
                        # Non-rate-limit error: break (don't retry)
                        break

        # ── Emergency Fallback: Groq ──
        if self.groq_available():
            print("[JarvisBrain] All Google keys exhausted. Falling back to Groq Cloud...")
            try:
                completion = self.groq_client.chat.completions.create(
                    model="qwen-2.5-coder-32b",
                    messages=[
                        {"role": "system", "content": system_instruction or "You are Jarvis, a coding assistant."},
                        {"role": "user", "content": prompt},
                    ],
                    temperature=0.6,
                    max_tokens=4096,
                )
                return completion.choices[0].message.content
            except Exception as e:
                return f"Both Google and Groq limits hit. Error: {e}"

        return "All free tier quotas exhausted. Please wait a short while for the reset window."

# ============================================================
# Backward-compatible globals for existing code
# ============================================================
jarvis_brain = JarvisAgentBrain()
USE_GEMINI = jarvis_brain.google_available()
USE_GROQ = jarvis_brain.groq_available()
GEMINI_MODEL_NAME = jarvis_brain.model_name
if USE_GEMINI:
    print(f"✓ Gemini API connected ({GEMINI_MODEL_NAME}) [via JarvisAgentBrain]")
if USE_GROQ:
    print("✓ Groq API connected - Ultra-fast mode enabled!")

# ============================================================
# GENERATIVE AI TOOLS & CAPABILITIES
# ============================================================

import subprocess
from datetime import datetime
from docx import Document
from docx.shared import Pt, Inches
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
import pyperclip
import psutil

# ============================================================
# PROFESSIONAL PRESENTATION ENGINE (v4.0) — 4-Step Design
# ============================================================
try:
    from professional_presentation import (
        create_professional_presentation,
        PROFESSIONAL_PRESENTATION_SYSTEM_PROMPT,
        auto_select_background,
        PROFESSIONAL_BACKGROUNDS,
    )
    PROFESSIONAL_MODE_AVAILABLE = True
    print("✓ Professional Presentation Engine v4.0 loaded (4-Step Design Process)")
except ImportError as e:
    print(f"  ⚠ Professional Presentation Engine not available: {e}")
    PROFESSIONAL_MODE_AVAILABLE = False
    PROFESSIONAL_PRESENTATION_SYSTEM_PROMPT = ""
    def auto_select_background(topic): return None
    PROFESSIONAL_BACKGROUNDS = {}

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "generated_files")
os.makedirs(OUTPUT_DIR, exist_ok=True)


# ============================================================
# FIX #3 — PROFESSIONAL PROMPT BUILDER
# Upgrades content generation to produce deeply structured,
# comprehensive, and professionally worded output.
# ============================================================

PROFESSIONAL_CONTENT_SYSTEM_PROMPT = """
You are an expert content writer and presentation specialist.
When generating content for documents or presentations, follow these rules strictly:

CONTENT QUALITY RULES:
1. Each bullet point must be a COMPLETE, INFORMATIVE sentence (minimum 15 words).
2. Use professional, authoritative language — no vague filler phrases.
3. Provide specific facts, statistics, or examples wherever possible.
4. Structure content logically: define → explain → imply → conclude.
5. For presentations: every slide must have 4-6 substantial bullet points.
6. For documents: write in full paragraphs (5-8 sentences each section).
7. Avoid one-word bullets like "Overview" or "Introduction" — always expand.
8. Use domain-specific vocabulary appropriate to the topic.
9. When generating slide content, think like a subject-matter expert giving a keynote.
10. Always generate AT LEAST 10 content slides for presentations.

SLIDE STRUCTURE TEMPLATE (follow this for every slide):
- Opening statement that defines the concept clearly
- 2-3 explanatory points with real-world context
- 1 example or case study reference
- 1 implication or takeaway point

OUTPUT FORMAT: Always respond with valid JSON only. No text outside JSON.
"""


def create_word_document(title, content, filename=None, show_live=True):
    """Create a Word document with the given title and content."""
    if filename is None:
        filename = f"document_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
    if not filename.endswith(".docx"):
        filename += ".docx"
    filepath = os.path.join(OUTPUT_DIR, filename)

    if show_live:
        try:
            import win32com.client as win32
            word = win32.Dispatch("Word.Application")
            word.Visible = True
            doc = word.Documents.Add()
            selection = word.Selection

            # Title
            selection.Font.Size = 18
            selection.Font.Bold = True
            selection.ParagraphFormat.Alignment = 1
            selection.TypeText(title)
            selection.TypeParagraph()
            selection.TypeParagraph()

            # Reset to body style
            selection.Font.Size = 12
            selection.Font.Bold = False
            selection.ParagraphFormat.Alignment = 0

            paragraphs = content.split("\n") if "\n" in content else [content]
            for para in paragraphs:
                if para.strip():
                    selection.TypeText(para.strip())
                    selection.TypeParagraph()
                    selection.TypeParagraph()  # Extra spacing between paragraphs

            doc.SaveAs2(filepath)
            word.Activate()

            # FIX #2 — Auto open document after creation
            os.startfile(filepath)
            return filepath

        except ImportError:
            print("[Info] pywin32 not installed, falling back to silent creation")
        except Exception as e:
            print(f"[Warning] Live creation failed: {e}, falling back to silent creation")

    # Silent fallback
    doc = Document()
    doc.add_heading(title, 0)
    if isinstance(content, list):
        for para in content:
            doc.add_paragraph(para)
    else:
        for para in content.split("\n"):
            if para.strip():
                doc.add_paragraph(para.strip())
    doc.save(filepath)

    # FIX #2 — Auto open even in silent mode
    try:
        os.startfile(filepath)
    except Exception as e:
        print(f"[Auto-open Warning]: {e}")

    return filepath


def download_image_from_pexels(query, save_path, slide_number=0):
    """Download a relevant image from Pexels based on query."""
    try:
        from pexels_integration import (
            download_background_image,
            get_color_palette_from_image,
        )
        downloaded_path = download_background_image(query, slide_number)
        if downloaded_path:
            import shutil
            shutil.copy2(downloaded_path, save_path)
            palette = get_color_palette_from_image(save_path)
            return {"path": save_path, "palette": palette, "query": query}
        return None
    except Exception as e:
        print(f"[Pexels Image Download Error]: {e}")
        return None


def create_presentation(title, slides_content, filename=None, show_live=True,
                        include_images=True, creative_mode=True):
    """
    Create a PowerPoint presentation.
    FIX #2: Auto-opens PowerPoint after generation.
    FIX #3: Uses professional content standards.
    """
    if creative_mode:
        try:
            from creative_presentation import create_creative_presentation
            print("\n  🎨 Using Creative Presentation Mode with Pexels backgrounds...")
            filepath = create_creative_presentation(title, slides_content, filename, show_progress=True)

            # FIX #2 — Auto-open PowerPoint immediately after generation
            _auto_open_file(filepath)
            return filepath
        except Exception as e:
            print(f"  ⚠ Creative mode failed: {e}, falling back to standard mode")

    if filename is None:
        filename = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    if not filename.endswith(".pptx"):
        filename += ".pptx"
    filepath = os.path.join(OUTPUT_DIR, filename)

    if show_live:
        try:
            import win32com.client as win32
            import pythoncom

            pythoncom.CoInitialize()
            print("\n  🎨 Opening PowerPoint for creative design...")
            ppt = win32.Dispatch("PowerPoint.Application")
            ppt.Visible = True
            ppt.WindowState = 1
            ppt.Activate()
            time.sleep(2)

            prs = ppt.Presentations.Add()
            prs.Windows(1).WindowState = 1

            print(f"\n  📝 Creating presentation: {title}")
            print(f"  📊 Total slides to create: {len(slides_content) + 1}")

            # Set to widescreen (13.333 x 7.5 inches)
            prs.PageSetup.SlideWidth = 13.333 * 72
            prs.PageSetup.SlideHeight = 7.5 * 72

            colors = [
                0x1A3A5C, 0x0D5E40, 0x6B1A1A, 0x2D1B5E,
                0x8B4500, 0x005858, 0x58005E, 0x1C3A1C,
                0x8B2222, 0x1B4B3A,
            ]

            # Title slide
            title_slide = prs.Slides.Add(1, 1)
            ppt.ActiveWindow.View.GotoSlide(1)
            title_shape = title_slide.Shapes(1).TextFrame.TextRange
            title_shape.Text = title
            title_shape.Font.Size = 54
            title_shape.Font.Bold = True
            title_shape.Font.Color.RGB = 0x1A3A5C
            title_shape.Font.Shadow = True

            subtitle_shape = title_slide.Shapes(2).TextFrame.TextRange
            subtitle_shape.Text = (
                f"A Comprehensive Professional Overview\n"
                f"Prepared by AI Assistant | {datetime.now().strftime('%B %Y')}"
            )
            subtitle_shape.Font.Size = 24
            subtitle_shape.Font.Italic = True
            time.sleep(1.5)

            for i, slide_data in enumerate(slides_content, 2):
                slide_title = slide_data.get("title", "Untitled")
                print(f"\n     🎨 Designing Slide {i - 1}: {slide_title}")

                layout = 2
                slide = prs.Slides.Add(i, layout)
                ppt.ActiveWindow.View.GotoSlide(i)
                time.sleep(0.8)

                title_text = slide.Shapes(1).TextFrame.TextRange
                title_text.Text = slide_title
                title_text.Font.Size = 36
                title_text.Font.Bold = True
                color_idx = (i - 2) % len(colors)
                title_text.Font.Color.RGB = colors[color_idx]
                title_text.Font.Shadow = True

                has_image = False
                img_path = None
                if include_images and "image_query" in slide_data:
                    try:
                        temp_path = os.path.join(OUTPUT_DIR, f"slide_{i}_image.jpg")
                        result = download_image_from_pexels(slide_data["image_query"], temp_path, i)
                        if result and os.path.exists(result["path"]):
                            has_image = True
                            img_path = result["path"]
                    except Exception as e:
                        print(f"        ⚠ Image skipped: {e}")

                if "content" in slide_data and len(slide.Shapes) > 1:
                    text_shape = slide.Shapes(2)
                    # Force strict word wrap to respect bounding box
                    text_shape.TextFrame.WordWrap = -1  # msoTrue

                    if has_image:
                        # Multi-column math: text on left, image on right
                        # Left col: 0.5" to 6.5" (Width: 6" = 432 pts)
                        text_shape.Left = 36
                        text_shape.Width = 432
                    else:
                        # Full width text
                        text_shape.Left = 36
                        text_shape.Width = 888

                    content_text = text_shape.TextFrame.TextRange
                    content_text.Text = slide_data["content"]
                    content_text.Font.Size = 18
                    content_text.ParagraphFormat.SpaceAfter = 14
                    content_text.ParagraphFormat.SpaceBefore = 4
                    try:
                        content_text.ParagraphFormat.Bullet.Visible = True
                    except Exception:
                        pass
                    time.sleep(1)

                if has_image:
                    try:
                        # Image strictly on the right half (Left: 6.8" = 490 pts, Width: 6" = 432 pts)
                        slide.Shapes.AddPicture(
                            img_path,
                            LinkToFile=False,
                            SaveWithDocument=True,
                            Left=490, Top=130, Width=432, Height=380,
                        )
                    except Exception as e:
                        print(f"        ⚠ Image layout failed: {e}")

                time.sleep(0.8)

            prs.SaveAs(filepath)
            print(f"\n  ✅ Presentation complete! {len(slides_content) + 1} slides created.")
            ppt.ActiveWindow.View.GotoSlide(1)
            ppt.Activate()
            ppt.WindowState = 1
            pythoncom.CoUninitialize()

            # FIX #2 — Auto-open PowerPoint instantly
            _auto_open_file(filepath)
            return filepath

        except ImportError:
            print("[Info] pywin32 not installed, falling back to silent creation")
        except Exception as e:
            print(f"[Warning] Live creation failed: {e}")
            import traceback
            traceback.print_exc()

    # Silent fallback
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    tbox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(2.5), PptxInches(12.333), PptxInches(2.0))
    tf = tbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = PptxPt(44)
    p.font.bold = True
    
    # Subtitle
    sbox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(4.5), PptxInches(12.333), PptxInches(1.0))
    stf = sbox.text_frame
    stf.word_wrap = True
    sp = stf.paragraphs[0]
    sp.text = "A Comprehensive Professional Overview\n" + datetime.now().strftime("%B %Y")
    sp.font.size = PptxPt(24)

    for i, slide_data in enumerate(slides_content):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
        
        # Add title
        tb = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.5), PptxInches(12.333), PptxInches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", "")
        p.font.size = PptxPt(36)
        p.font.bold = True
        
        has_image = False
        img_path = None
        if include_images and "image_query" in slide_data:
            try:
                temp_path = os.path.join(OUTPUT_DIR, f"slide_fallback_{i}_image.jpg")
                result = download_image_from_pexels(slide_data["image_query"], temp_path, i)
                if result and os.path.exists(result["path"]):
                    has_image = True
                    img_path = result["path"]
            except Exception:
                pass

        if has_image:
            # Multi-column math: Text on left, Image on right
            content_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(1.6), PptxInches(6.0), PptxInches(5.5))
            ctf = content_box.text_frame
            ctf.word_wrap = True
            if "content" in slide_data:
                ctf.text = slide_data["content"]
            
            # Image container
            slide.shapes.add_picture(img_path, PptxInches(6.8), PptxInches(1.6), width=PptxInches(6.0))
        else:
            # Full width
            content_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(1.6), PptxInches(12.0), PptxInches(5.5))
            ctf = content_box.text_frame
            ctf.word_wrap = True
            if "content" in slide_data:
                ctf.text = slide_data["content"]

    prs.save(filepath)

    # FIX #2 — Auto-open even in silent mode
    _auto_open_file(filepath)
    return filepath


def _auto_open_file(filepath: str):
    """
    FIX #2 — Cross-platform auto-open for generated files.
    Opens the file immediately using the OS default application.
    """
    try:
        if os.name == "nt":  # Windows
            os.startfile(filepath)
        elif os.uname().sysname == "Darwin":  # macOS
            subprocess.Popen(["open", filepath])
        else:  # Linux
            subprocess.Popen(["xdg-open", filepath])
        print(f"  🚀 Auto-opened: {os.path.basename(filepath)}")
    except Exception as e:
        print(f"  [Auto-open failed]: {e} — Please open manually: {filepath}")


def create_resume(name, contact_info, experience, education, skills,
                  filename=None, show_live=True):
    """Create a professional resume."""
    if filename is None:
        filename = f"resume_{name.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
    if not filename.endswith(".docx"):
        filename += ".docx"
    filepath = os.path.join(OUTPUT_DIR, filename)

    if show_live:
        try:
            import win32com.client as win32
            word = win32.Dispatch("Word.Application")
            word.Visible = True
            doc = word.Documents.Add()
            selection = word.Selection

            selection.Font.Size = 20
            selection.Font.Bold = True
            selection.ParagraphFormat.Alignment = 1
            selection.TypeText(name)
            selection.TypeParagraph()

            selection.Font.Size = 11
            selection.Font.Bold = False
            selection.TypeText(contact_info)
            selection.TypeParagraph()
            selection.TypeParagraph()

            if experience:
                selection.Font.Size = 14
                selection.Font.Bold = True
                selection.ParagraphFormat.Alignment = 0
                selection.TypeText("PROFESSIONAL EXPERIENCE")
                selection.TypeParagraph()
                selection.Font.Size = 12
                selection.Font.Bold = False
                for exp in experience:
                    selection.Font.Bold = True
                    selection.TypeText(exp.get("title", ""))
                    selection.Font.Bold = False
                    selection.TypeText(f" — {exp.get('company', '')}")
                    if exp.get("duration"):
                        selection.TypeText(f"  |  {exp.get('duration', '')}")
                    selection.TypeParagraph()
                    if exp.get("description"):
                        selection.TypeText(exp["description"])
                        selection.TypeParagraph()
                    selection.TypeParagraph()

            if education:
                selection.Font.Size = 14
                selection.Font.Bold = True
                selection.TypeText("EDUCATION")
                selection.TypeParagraph()
                selection.Font.Size = 12
                selection.Font.Bold = False
                for edu in education:
                    selection.Font.Bold = True
                    selection.TypeText(edu.get("degree", ""))
                    selection.Font.Bold = False
                    selection.TypeText(f" — {edu.get('institution', '')}")
                    if edu.get("year"):
                        selection.TypeText(f"  ({edu.get('year', '')})")
                    selection.TypeParagraph()
                selection.TypeParagraph()

            if skills:
                selection.Font.Size = 14
                selection.Font.Bold = True
                selection.TypeText("CORE SKILLS")
                selection.TypeParagraph()
                selection.Font.Size = 12
                selection.Font.Bold = False
                selection.TypeText(" • ".join(skills))

            doc.SaveAs2(filepath)
            word.Activate()
            _auto_open_file(filepath)
            return filepath

        except ImportError:
            print("[Info] pywin32 not installed, falling back to silent creation")
        except Exception as e:
            print(f"[Warning] Live creation failed: {e}")

    # Silent fallback
    doc = Document()
    doc.add_heading(name, 0)
    doc.add_paragraph(contact_info)
    if experience:
        doc.add_heading("Professional Experience", level=1)
        for exp in experience:
            p = doc.add_paragraph()
            p.add_run(exp.get("title", "")).bold = True
            p.add_run(f" — {exp.get('company', '')}")
            if exp.get("duration"):
                p.add_run(f"  |  {exp.get('duration', '')}")
            if exp.get("description"):
                doc.add_paragraph(exp["description"])
    if education:
        doc.add_heading("Education", level=1)
        for edu in education:
            p = doc.add_paragraph()
            p.add_run(edu.get("degree", "")).bold = True
            p.add_run(f" — {edu.get('institution', '')}")
            if edu.get("year"):
                p.add_run(f"  ({edu.get('year', '')})")
    if skills:
        doc.add_heading("Core Skills", level=1)
        doc.add_paragraph(" • ".join(skills))
    doc.save(filepath)
    _auto_open_file(filepath)
    return filepath


# ============================================================
# WEBSITE / WEB APP DATABASE — Opens in default browser
# ============================================================
WEBSITE_MAP = {
    "youtube": "https://www.youtube.com",
    "yt": "https://www.youtube.com",
    "youtube music": "https://music.youtube.com",
    "youtube studio": "https://studio.youtube.com",
    "google": "https://www.google.com",
    "gmail": "https://mail.google.com",
    "mail": "https://mail.google.com",
    "email": "https://mail.google.com",
    "google drive": "https://drive.google.com",
    "drive": "https://drive.google.com",
    "google docs": "https://docs.google.com",
    "docs": "https://docs.google.com",
    "google sheets": "https://sheets.google.com",
    "sheets": "https://sheets.google.com",
    "google slides": "https://slides.google.com",
    "slides": "https://slides.google.com",
    "google meet": "https://meet.google.com",
    "meet": "https://meet.google.com",
    "google calendar": "https://calendar.google.com",
    "calendar": "https://calendar.google.com",
    "google maps": "https://maps.google.com",
    "maps": "https://maps.google.com",
    "github": "https://github.com",
    "git hub": "https://github.com",
    "stack overflow": "https://stackoverflow.com",
    "stackoverflow": "https://stackoverflow.com",
    "facebook": "https://www.facebook.com",
    "fb": "https://www.facebook.com",
    "instagram": "https://www.instagram.com",
    "insta": "https://www.instagram.com",
    "whatsapp": "https://web.whatsapp.com",
    "whats app": "https://web.whatsapp.com",
    "wa": "https://web.whatsapp.com",
    "twitter": "https://www.twitter.com",
    "x": "https://www.twitter.com",
    "linkedin": "https://www.linkedin.com",
    "reddit": "https://www.reddit.com",
    "amazon": "https://www.amazon.in",
    "flipkart": "https://www.flipkart.com",
    "netflix": "https://www.netflix.com",
    "hotstar": "https://www.hotstar.com",
    "prime video": "https://www.primevideo.com",
    "amazon prime": "https://www.primevideo.com",
    "spotify web": "https://open.spotify.com",
    "chatgpt": "https://chat.openai.com",
    "chat gpt": "https://chat.openai.com",
    "openai": "https://chat.openai.com",
    "bard": "https://bard.google.com",
    "gemini": "https://gemini.google.com",
    "google ai": "https://gemini.google.com",
    "claude": "https://claude.ai",
    "perplexity": "https://www.perplexity.ai",
    "canva": "https://www.canva.com",
    "zoom web": "https://zoom.us",
    "teams web": "https://teams.microsoft.com",
    "slack web": "https://slack.com",
    "notion": "https://www.notion.so",
    "medium": "https://medium.com",
    "dev.to": "https://dev.to",
    "wikipedia": "https://en.wikipedia.org",
    "wiki": "https://en.wikipedia.org",
    "quora": "https://www.quora.com",
    "pinterest": "https://in.pinterest.com",
    "telegram web": "https://web.telegram.org",
    "discord web": "https://discord.com/app",
    "trello": "https://trello.com",
    "jira": "https://www.atlassian.com/software/jira",
    "bitbucket": "https://bitbucket.org",
    "gitlab": "https://gitlab.com",
    "npm": "https://www.npmjs.com",
    "pypi": "https://pypi.org",
    "pythonanywhere": "https://www.pythonanywhere.com",
    "replit": "https://replit.com",
    "code pen": "https://codepen.io",
    "codepen": "https://codepen.io",
    "figma": "https://www.figma.com",
    "dribbble": "https://dribbble.com",
    "behance": "https://www.behance.net",
    "unsplash": "https://unsplash.com",
    "pexels": "https://www.pexels.com",
    "pixabay": "https://pixabay.com",
    "spotify": "https://open.spotify.com",
}


# ============================================================
# SYSTEM / DESKTOP APPLICATION MAP (exe names + aliases)
# ============================================================
# Keys are lowercase aliases users might say.
# Values are what we pass to os.startfile() — can be a .exe name,
# a protocol URI (ms-settings:), or a direct path.
# ============================================================
APP_MAP = {
    # ── Browsers ──
    "chrome": "chrome",
    "google chrome": "chrome",
    "browser": "chrome",
    "edge": "msedge",
    "microsoft edge": "msedge",
    "firefox": "firefox",
    "mozilla firefox": "firefox",
    "brave": "brave",
    "opera": "opera",
    "browser brave": "brave",
    "browser chrome": "chrome",
    "browser firefox": "firefox",

    # ── Microsoft Office ──
    "word": "WINWORD.EXE",
    "microsoft word": "WINWORD.EXE",
    "ms word": "WINWORD.EXE",
    "excel": "EXCEL.EXE",
    "microsoft excel": "EXCEL.EXE",
    "ms excel": "EXCEL.EXE",
    "powerpoint": "POWERPNT.EXE",
    "ppt": "POWERPNT.EXE",
    "microsoft powerpoint": "POWERPNT.EXE",
    "ms powerpoint": "POWERPNT.EXE",
    "outlook": "OUTLOOK.EXE",
    "microsoft outlook": "OUTLOOK.EXE",
    "ms outlook": "OUTLOOK.EXE",
    "onenote": "ONENOTE.EXE",
    "microsoft onenote": "ONENOTE.EXE",
    "access": "MSACCESS.EXE",
    "microsoft access": "MSACCESS.EXE",
    "publisher": "MSPUB.EXE",
    "microsoft publisher": "MSPUB.EXE",
    "teams": "teams",
    "microsoft teams": "teams",
    "ms teams": "teams",

    # ── Windows System Tools ──
    "notepad": "notepad",
    "notepad++": "notepad++",
    "calculator": "calc",
    "calc": "calc",
    "paint": "mspaint",
    "microsoft paint": "mspaint",
    "file explorer": "explorer",
    "explorer": "explorer",
    "windows explorer": "explorer",
    "cmd": "cmd",
    "command prompt": "cmd",
    "terminal": "cmd",
    "powershell": "powershell",
    "windows powershell": "powershell",
    "task manager": "taskmgr",
    "taskmgr": "taskmgr",
    "control panel": "control",
    "settings": "ms-settings:",
    "windows settings": "ms-settings:",
    "device manager": "devmgmt.msc",
    "disk management": "diskmgmt.msc",
    "disk cleanup": "cleanmgr",
    "registry editor": "regedit",
    "regedit": "regedit",
    "system information": "msinfo32",
    "resource monitor": "resmon",
    "performance monitor": "perfmon",
    "services": "services.msc",
    "task scheduler": "taskschd.msc",
    "event viewer": "eventvwr.msc",
    "snipping tool": "SnippingTool",
    "snip and sketch": "SnippingTool",
    "screenshot tool": "SnippingTool",
    "magnifier": "Magnify",
    "narrator": "Narrator",
    "on screen keyboard": "osk",
    "osk": "osk",
    "character map": "charmap",
    "command": "cmd",

    # ── Windows Store / Modern Apps ──
    "camera": "microsoft.windows.camera:",
    "windows camera": "microsoft.windows.camera:",
    "clock": "ms-clock:",
    "alarms": "ms-clock:",
    "alarm": "ms-clock:",
    "calculator (modern)": "calculator:",
    "calendar (modern)": "outlookcal:",
    "mail (modern)": "outlookmail:",
    "maps (modern)": "bingmaps:",
    "news": "msnnews:",
    "weather": "bingweather:",
    "xbox": "xbox:",
    "xbox app": "xbox:",
    "store": "ms-windows-store:",
    "microsoft store": "ms-windows-store:",
    "windows store": "ms-windows-store:",
    "photos": "ms-photos:",
    "windows photos": "ms-photos:",
    "music": "ms-wvvmusic:",
    "groove music": "ms-wvvmusic:",
    "video": "ms-wvvideo:",
    "movies & tv": "ms-wvvideo:",
    "tips": "ms-get-started:",
    "feedback hub": "feedback-hub:",
    "solitare": "ms-solitaire:",
    "microsoft solitare": "ms-solitaire:",
    "whiteboard": "whiteboard:",
    "microsoft whiteboard": "whiteboard:",

    # ── Development Tools ──
    "vscode": "code",
    "vs code": "code",
    "visual studio code": "code",
    "code": "code",
    "visual studio": "devenv",
    "vs": "devenv",
    "pycharm": "pycharm",
    "jetbrains pycharm": "pycharm",
    "intellij": "idea",
    "intellij idea": "idea",
    "webstorm": "webstorm",
    "phpstorm": "phpstorm",
    "android studio": "studio64",
    "eclipse": "eclipse",
    "sublime text": "sublime_text",
    "atom": "atom",
    "git bash": "git-bash",
    "git": "git-bash",
    "postman": "postman",
    "docker": "docker",
    "docker desktop": "docker",
    "node.js": "node",
    "python (idle)": "idle",
    "thonny": "thonny",
    "jupyter": "jupyter-notebook",
    "jupyter notebook": "jupyter-notebook",
    "anaconda": "anaconda-navigator",
    "anaconda navigator": "anaconda-navigator",
    "vmware": "vmware",
    "virtualbox": "virtualbox",
    "putty": "putty",
    "winscp": "winscp",
    "filezilla": "filezilla",
    "mysql workbench": "mysqlworkbench",
    "sql server": "ssms",
    "ssms": "ssms",
    "azure data studio": "azuredatastudio",

    # ── Communication & Social ──
    "discord": "discord",
    "slack": "slack",
    "zoom": "zoom",
    "zoom meeting": "zoom",
    "zoom client": "zoom",
    "skype": "skype",
    "telegram": "telegram",
    "whatsapp desktop": "whatsapp",
    "whatsapp": "whatsapp",
    "signal": "signal",
    "thunderbird": "thunderbird",
    "mozilla thunderbird": "thunderbird",

    # ── Media & Entertainment ──
    "vlc": "vlc",
    "vlc media player": "vlc",
    "media player": "wmplayer",
    "windows media player": "wmplayer",
    "spotify": "spotify",
    "itunes": "itunes",
    "kodi": "kodi",
    "plex": "plex",
    "foobar2000": "foobar2000",
    "audacity": "audacity",
    "obs": "obs64",
    "obs studio": "obs64",
    "streamlabs obs": "streamlabs-obs",
    "steam": "steam",
    "epic games": "epicgameslauncher",
    "epic games launcher": "epicgameslauncher",

    # ── Graphics & Design ──
    "photoshop": "photoshop",
    "adobe photoshop": "photoshop",
    "illustrator": "illustrator",
    "adobe illustrator": "illustrator",
    "premiere pro": "adobepremierepro",
    "adobe premiere pro": "adobepremierepro",
    "after effects": "afterfx",
    "adobe after effects": "afterfx",
    "lightroom": "lightroom",
    "adobe lightroom": "lightroom",
    "adobe acrobat": "acrobat",
    "acrobat": "acrobat",
    "acrobat reader": "acrobat",
    "pdf reader": "acrobat",
    "gimp": "gimp",
    "blender": "blender",
    "coreldraw": "coreldraw",
    "inkscape": "inkscape",
    "figma (desktop)": "figma",
    "canva (desktop)": "canva",

    # ── Utilities ──
    "winrar": "winrar",
    "win rar": "winrar",
    "7-zip": "7zfm",
    "7zip": "7zfm",
    "7 zip": "7zfm",
    "cpu-z": "cpuz",
    "gpu-z": "gpuz",
    "hwmonitor": "hwmonitor",
    "hw monitor": "hwmonitor",
    "speccy": "speccy",
    "ccleaner": "ccleaner",
    "malwarebytes": "mbam",
    "antivirus": "mbam",
    "defender": "WindowsDefender",
    "windows defender": "WindowsDefender",
    "one drive": "onedrive",
    "onedrive": "onedrive",
    "dropbox": "dropbox",
    "google drive (desktop)": "googledrivesync",
    "evernote": "evernote",
    "lastpass": "lastpass",
    "1password": "1password",
}


def _find_app_in_registry(app_name_lower: str) -> str | None:
    """
    Search Windows Registry for installed applications by name.
    Checks both 64-bit and 32-bit registry locations.
    Returns the full path to the executable if found, None otherwise.
    """
    import winreg
    
    registry_paths = [
        (winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths"),
        (winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\WOW6432Node\Microsoft\Windows\CurrentVersion\App Paths"),
        (winreg.HKEY_CURRENT_USER, r"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths"),
    ]
    
    # First try: search by exact name + .exe
    exe_name = app_name_lower.strip()
    if not exe_name.endswith('.exe'):
        exe_name += '.exe'
    
    for hive, path in registry_paths:
        try:
            with winreg.OpenKey(hive, path) as key:
                try:
                    with winreg.OpenKey(key, exe_name) as subkey:
                        value, _ = winreg.QueryValueEx(subkey, "")
                        if value and os.path.exists(value):
                            return value
                except FileNotFoundError:
                    pass
        except FileNotFoundError:
            continue
    
    # Second try: enumerate all keys and do a partial name match
    for hive, path in registry_paths:
        try:
            with winreg.OpenKey(hive, path) as key:
                i = 0
                while True:
                    try:
                        subkey_name = winreg.EnumKey(key, i)
                        i += 1
                        # Check if the app name contains our search term or vice versa
                        base_name = subkey_name.lower().replace('.exe', '').replace('_', ' ')
                        if (app_name_lower in base_name or base_name in app_name_lower):
                            try:
                                with winreg.OpenKey(key, subkey_name) as subkey:
                                    value, _ = winreg.QueryValueEx(subkey, "")
                                    if value and os.path.exists(value):
                                        return value
                            except (FileNotFoundError, OSError):
                                pass
                    except FileNotFoundError:
                        break
        except FileNotFoundError:
            continue
    
    return None


def _find_app_in_start_menu(app_name_lower: str) -> str | None:
    """
    Search Windows Start Menu shortcuts (.lnk files) for the app.
    Checks both All Users and Current User start menu directories.
    Returns the full path to the .lnk file if found.
    """
    start_menu_paths = [
        os.path.join(os.environ.get("PROGRAMDATA", "C:\\ProgramData"), "Microsoft", "Windows", "Start Menu", "Programs"),
        os.path.join(os.environ.get("APPDATA", ""), "Microsoft", "Windows", "Start Menu", "Programs"),
    ]
    
    for base_path in start_menu_paths:
        if not os.path.exists(base_path):
            continue
        for root, dirs, files in os.walk(base_path):
            for file in files:
                if file.endswith('.lnk'):
                    file_lower = file.lower().replace('.lnk', '')
                    # Check if app name matches the shortcut name
                    if (app_name_lower in file_lower or file_lower in app_name_lower):
                        return os.path.join(root, file)
                    # Also check if folder name matches (e.g., "Spotify" folder)
                    folder_name = os.path.basename(root).lower()
                    if app_name_lower in folder_name or folder_name in app_name_lower:
                        return os.path.join(root, file)
        # Also check folder names at the top level
        for d in dirs:
            if app_name_lower in d.lower() or d.lower() in app_name_lower:
                folder_path = os.path.join(base_path, d)
                for f in os.listdir(folder_path):
                    if f.endswith('.lnk'):
                        return os.path.join(folder_path, f)
    
    return None


def _find_app_in_program_files(app_name_lower: str) -> str | None:
    """
    Search Program Files directories for executables matching the app name.
    """
    program_dirs = [
        os.environ.get("ProgramFiles", "C:\\Program Files"),
        os.environ.get("ProgramFiles(x86)", "C:\\Program Files (x86)"),
        os.environ.get("LOCALAPPDATA", ""),
        os.environ.get("APPDATA", ""),
    ]
    
    search_exe = app_name_lower.strip()
    if not search_exe.endswith('.exe'):
        search_exe += '.exe'
    
    for base_dir in program_dirs:
        if not base_dir or not os.path.exists(base_dir):
            continue
        try:
            for root, dirs, files in os.walk(base_dir, topdown=True):
                # Limit search depth to avoid extreme delays
                depth = root.replace(base_dir, '').count(os.sep)
                if depth > 4:
                    dirs.clear()  # Don't go deeper
                    continue
                for file in files:
                    file_lower = file.lower()
                    if file_lower == search_exe:
                        full_path = os.path.join(root, file)
                        if os.path.exists(full_path):
                            return full_path
                    # Also check if the folder name matches (common pattern)
                    folder_name = os.path.basename(root).lower()
                    if (app_name_lower in folder_name or folder_name in app_name_lower) and file_lower.endswith('.exe'):
                        # Pick the first reasonable exe (not uninstall, not setup)
                        if 'unins' not in file_lower and 'setup' not in file_lower and 'install' not in file_lower:
                            full_path = os.path.join(root, file)
                            if os.path.exists(full_path):
                                return full_path
        except (PermissionError, OSError):
            continue
    
    return None


def _find_app_on_path(app_name_lower: str) -> str | None:
    """
    Search the system PATH for the executable (using `where` command on Windows).
    """
    try:
        exe_name = app_name_lower.strip()
        if not exe_name.endswith('.exe'):
            exe_name += '.exe'
        result = subprocess.run(
            ["where", exe_name],
            capture_output=True, text=True, timeout=5
        )
        if result.returncode == 0 and result.stdout.strip():
            paths = [p.strip() for p in result.stdout.strip().split('\n') if p.strip()]
            for p in paths:
                if os.path.exists(p):
                    return p
    except Exception:
        pass
    return None


def open_application(app_name):
    """
    Open an application or website by name.
    
    First checks if the name matches a known website (opens in default browser).
    Then checks the app map for known applications.
    Then searches Windows Registry, Start Menu, Program Files, and PATH.
    
    Supports:
    - Web apps: "youtube", "gmail", "github", "facebook", "instagram", etc.
    - System apps: "notepad", "calculator", "cmd", etc.
    - Office apps: "word", "excel", "powerpoint", etc.
    - Installed apps: "spotify", "discord", "vscode", "steam", etc.
    - Any app found in Windows Registry or Start Menu
    """
    try:
        app_lower = app_name.strip().lower()
        
        # ── STEP 1: Check if it's a website/web app ──
        # Match against known website names
        for key, url in WEBSITE_MAP.items():
            if key in app_lower or app_lower in key:
                import webbrowser
                webbrowser.open(url)
                return f"Opened {app_name} in your default browser"
        
        # If it contains common web patterns, try as URL
        if any(ext in app_lower for ext in ['.com', '.org', '.net', '.io', '.ai', '.app']):
            if not app_lower.startswith('http'):
                url = 'https://' + app_lower
            else:
                url = app_lower
            import webbrowser
            webbrowser.open(url)
            return f"Opened {app_name} in your default browser"
        
        # ── STEP 2: Check the app map ──
        for key, target in APP_MAP.items():
            if key in app_lower or app_lower in key:
                os.startfile(target)
                return f"Opened {app_name}"
        
        # ── STEP 3: Search Windows Registry ──
        registry_path = _find_app_in_registry(app_lower)
        if registry_path:
            os.startfile(registry_path)
            return f"Opened {app_name} (found in Windows Registry)"
        
        # ── STEP 4: Search Start Menu ──
        start_menu_path = _find_app_in_start_menu(app_lower)
        if start_menu_path:
            os.startfile(start_menu_path)
            return f"Opened {app_name} (found in Start Menu)"
        
        # ── STEP 5: Search Program Files ──
        program_files_path = _find_app_in_program_files(app_lower)
        if program_files_path:
            os.startfile(program_files_path)
            return f"Opened {app_name} (found in Program Files)"
        
        # ── STEP 6: Search PATH ──
        path_exe = _find_app_on_path(app_lower)
        if path_exe:
            os.startfile(path_exe)
            return f"Opened {app_name} (found on system PATH)"
        
        # ── STEP 7: Final fallback — try os.startfile directly ──
        # This may fail but we attempt it as a last resort
        try:
            os.startfile(app_lower)
            return f"Opened {app_name}"
        except Exception:
            pass
        
        return f"Could not find {app_name}. Please check the name and try again."
        
    except Exception as e:
        return f"Failed to open {app_name}: {str(e)}"


def get_system_info():
    """Get current system information including CPU, memory, disk, battery, and GPU."""
    try:
        info = {
            "cpu_usage": f"{psutil.cpu_percent(interval=1)}%",
            "cpu_count": f"{psutil.cpu_count()} cores",
            "memory": f"{psutil.virtual_memory().percent}%",
            "memory_total": f"{psutil.virtual_memory().total / (1024**3):.1f} GB",
            "memory_available": f"{psutil.virtual_memory().available / (1024**3):.1f} GB",
            "disk_usage": f"{psutil.disk_usage('C:').percent}%",
            "disk_total": f"{psutil.disk_usage('C:').total / (1024**3):.1f} GB",
            "disk_free": f"{psutil.disk_usage('C:').free / (1024**3):.1f} GB",
        }
        battery = psutil.sensors_battery()
        if battery:
            info["battery"] = f"{battery.percent}%"
            info["battery_plugged"] = battery.power_plugged
        else:
            info["battery"] = "No battery detected (desktop system)"
        
        # Try to get GPU info using Windows Management Instrumentation
        try:
            import subprocess
            # Try nvidia-smi first (NVIDIA GPUs)
            gpu_result = subprocess.run(
                ["nvidia-smi", "--query-gpu=name,utilization.gpu,memory.used,memory.total", "--format=csv,noheader,nounits"],
                capture_output=True, text=True, timeout=5
            )
            if gpu_result.returncode == 0 and gpu_result.stdout.strip():
                gpu_lines = gpu_result.stdout.strip().split('\n')
                gpus = []
                for line in gpu_lines:
                    if ',' in line:
                        parts = [p.strip() for p in line.split(',')]
                        if len(parts) >= 4:
                            gpus.append({
                                "name": parts[0],
                                "utilization": f"{parts[1]}%",
                                "memory_used": f"{parts[2]} MB",
                                "memory_total": f"{parts[3]} MB"
                            })
                if gpus:
                    info["gpu"] = gpus
            else:
                # Fallback: try wmic on Windows
                wmic_result = subprocess.run(
                    ["wmic", "path", "win32_videocontroller", "get", "name"],
                    capture_output=True, text=True, timeout=5
                )
                if wmic_result.returncode == 0:
                    lines = [l.strip() for l in wmic_result.stdout.split('\n') if l.strip()]
                    if len(lines) > 1:
                        gpu_name = lines[1]
                        info["gpu"] = [{"name": gpu_name}]
        except Exception:
            pass
        
        return info
    except Exception as e:
        return {"error": str(e), "message": "Could not retrieve system information"}


def write_to_clipboard(text):
    pyperclip.copy(text)
    return "Text copied to clipboard"


def read_clipboard():
    return pyperclip.paste()


def list_files(directory=None):
    if directory is None:
        directory = os.path.expanduser("~\\Desktop")
    try:
        files = os.listdir(directory)
        return {"directory": directory, "files": files[:20]}
    except Exception as e:
        return {"error": str(e)}


def create_file(filename, content, folder="generated_files"):
    filepath = os.path.join(OUTPUT_DIR, filename)
    with open(filepath, "w", encoding="utf-8") as f:
        f.write(content)
    return filepath


# ============================================================
# AVAILABLE TOOLS DEFINITION
# ============================================================

AVAILABLE_TOOLS = [
    {
        "type": "function",
        "function": {
            "name": "create_word_document",
            "description": "Create a detailed, professional Word document with title and comprehensive content",
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {"type": "string", "description": "Document title"},
                    "content": {"type": "string", "description": "Full document content — must be detailed, multi-paragraph, professional"},
                    "filename": {"type": "string", "description": "Optional filename ending in .docx"},
                },
                "required": ["title", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "create_presentation",
            "description": "Create a professional PowerPoint presentation with at least 10 slides, each with detailed content and images",
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {"type": "string", "description": "Presentation title"},
                    "slides_content": {
                        "type": "string",
                        "description": (
                            "JSON array of slide objects. Each object must have: "
                            "title (string), content (4-6 detailed bullet points separated by \\n), "
                            "image_query (search term for background image). "
                            "Minimum 10 slides required."
                        ),
                    },
                    "filename": {"type": "string", "description": "Optional filename ending in .pptx"},
                },
                "required": ["title", "slides_content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "create_resume",
            "description": "Create a professional resume document",
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {"type": "string"},
                    "contact_info": {"type": "string"},
                    "experience": {"type": "string", "description": "JSON array of experience objects"},
                    "education": {"type": "string", "description": "JSON array of education objects"},
                    "skills": {"type": "string", "description": "JSON array of skills"},
                },
                "required": ["name", "contact_info"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "open_application",
            "description": "Open an application on the computer",
            "parameters": {
                "type": "object",
                "properties": {"app_name": {"type": "string"}},
                "required": ["app_name"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "get_system_info",
            "description": "Get current system information (CPU, memory, battery)",
            "parameters": {"type": "object", "properties": {}},
        },
    },
    {
        "type": "function",
        "function": {
            "name": "write_to_clipboard",
            "description": "Write text to clipboard",
            "parameters": {
                "type": "object",
                "properties": {"text": {"type": "string"}},
                "required": ["text"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "create_file",
            "description": "Create a text file with content",
            "parameters": {
                "type": "object",
                "properties": {
                    "filename": {"type": "string"},
                    "content": {"type": "string"},
                },
                "required": ["filename", "content"],
            },
        },
    },
]

FUNCTION_MAP = {
    "create_word_document": create_word_document,
    "create_presentation": create_presentation,
    "create_resume": create_resume,
    "open_application": open_application,
    "get_system_info": get_system_info,
    "write_to_clipboard": write_to_clipboard,
    "create_file": create_file,
}


def execute_function(name, arguments):
    """Execute a function with the given arguments."""
    if name not in FUNCTION_MAP:
        return f"Unknown function: {name}"
    try:
        func = FUNCTION_MAP[name]
        parsed_args = {}
        # FIX: Handle None/null arguments from LLM JSON output
        if arguments is None:
            arguments = {}
        if arguments:
            for key, value in arguments.items():
                if isinstance(value, str) and value.strip().startswith(("[", "{")):
                    try:
                        parsed_args[key] = json.loads(value)
                    except Exception:
                        parsed_args[key] = value
                else:
                    parsed_args[key] = value
        result = func(**parsed_args)
        return f"Success: {result}"
    except Exception as e:
        return f"Error executing {name}: {str(e)}"


# ============================================================
# SPEECH RECOGNITION (voice input via sounddevice)
# Uses sounddevice as microphone backend instead of PyAudio
# ============================================================
import threading
import queue
import wave
import tempfile

try:
    import speech_recognition as sr
    SR_AVAILABLE = True
    print("✓ SpeechRecognition loaded — voice input enabled!")
except ImportError:
    SR_AVAILABLE = False
    print("  ⚠ SpeechRecognition not installed — voice input disabled")

# ── Custom Microphone class using sounddevice (no PyAudio needed) ──
class SoundDeviceMicrophone(sr.AudioSource):
    """A microphone class for SpeechRecognition that uses sounddevice instead of PyAudio."""
    
    def __init__(self, sample_rate=16000, channels=1, device=None, timeout=5):
        self.SAMPLE_RATE = sample_rate
        self.SAMPLE_WIDTH = 2  # 16-bit audio = 2 bytes per sample
        self.CHANNELS = channels
        self.device = device
        self.CHUNK = 1024
        self.timeout = timeout
        self.audio_queue = queue.Queue()
        self.stop_event = threading.Event()
        self.recording_thread = None

    def __enter__(self):
        self._sd_stream = sd.InputStream(
            samplerate=self.SAMPLE_RATE,
            channels=self.CHANNELS,
            device=self.device,
            dtype='int16',
            blocksize=self.CHUNK,
        )
        self._sd_stream.__enter__()
        return self

    def __exit__(self, exc_type, exc_val, exc_tb):
        self._sd_stream.__exit__(exc_type, exc_val, exc_tb)

    @property
    def stream(self):
        """SpeechRecognition calls source.stream.read() — return self so our read() is used."""
        return self

    def read(self, size):
        """Read audio data from the stream. Size is in bytes — convert to frames for sounddevice."""
        frames_needed = max(1, size // (self.SAMPLE_WIDTH * self.CHANNELS))
        frames, _ = self._sd_stream.read(frames_needed)
        return frames.tobytes()


# ============================================================
# NATURAL VOICE ACTIVITY DETECTION (VAD)
# Event-driven listening — waits for you to finish speaking naturally
# Uses WebRTC VAD + silence detection for natural conversation pacing
# ============================================================

class NaturalVoiceListener:
    """
    Listens to the microphone using a continuous stream buffer,
    detects speech vs silence in real-time, and only stops recording
    AFTER a natural pause in speech (1.5s of silence).
    
    This prevents the "too fast" problem where the system cuts you off
    mid-sentence because of arbitrary timeouts.
    """
    
    # WebRTC VAD mode 0 = aggressive (best for quiet rooms)
    # mode 1 = less aggressive (for noisier environments)
    # mode 2 = even less aggressive
    # mode 3 = least aggressive (= more sensitive, detects softer speech)
    VAD_MODE = 1  
    
    # Silence threshold: how many seconds of quiet before we decide speech is done
    SILENCE_DURATION_SECONDS = 1.5
    
    # Minimum speech duration to accept (avoids coughs/background noise)
    MIN_SPEECH_SECONDS = 0.5
    
    # Max recording duration (safety limit — 2 minutes)
    MAX_RECORD_SECONDS = 120
    
    def __init__(self, sample_rate=16000):
        self.sample_rate = sample_rate
        self.channels = 1
        self.sample_width = 2  # 16-bit
        self.chunk_size = 480  # 30ms frames at 16kHz (WebRTC VAD standard)
        
    def record_until_silence(self) -> bytes:
        """
        Record audio from microphone continuously.
        Returns raw PCM audio bytes of the full utterance.
        
        Algorithm:
        1. Buffer audio in 30ms frames
        2. Measure RMS energy per frame
        3. Track speech/silence state transitions
        4. After speech detected, wait for SILENCE_DURATION_SECONDS of quiet
        5. Return the complete audio buffer
        """
        audio_buffer = bytearray()
        is_speaking = False
        silence_counter = 0
        speech_counter = 0
        
        # Calculate frame sizes
        frames_for_silence = int(self.SILENCE_DURATION_SECONDS * self.sample_rate / self.chunk_size)
        frames_for_min_speech = int(self.MIN_SPEECH_SECONDS * self.sample_rate / self.chunk_size)
        max_frames = int(self.MAX_RECORD_SECONDS * self.sample_rate / self.chunk_size)
        
        # Energy thresholds (calibrated for 16-bit audio)
        speech_threshold = 500    # RMS above this = speech
        silence_threshold = 300   # RMS below this = silence
        
        # Calibration phase
        print("  🔇 [Calibrating ambient noise...]", end="", flush=True)
        noise_samples = []
        calibration_frames = int(0.5 * self.sample_rate / self.chunk_size)  # 0.5 seconds
        
        with sd.InputStream(
            samplerate=self.sample_rate,
            channels=self.channels,
            dtype='int16',
            blocksize=self.chunk_size,
        ) as stream:
            for _ in range(calibration_frames):
                frame, _ = stream.read(self.chunk_size)
                rms = np.sqrt(np.mean(frame.astype(np.float32)**2))
                noise_samples.append(rms)
            
            if noise_samples:
                ambient_rms = np.median(noise_samples)
                # Set threshold: 3x ambient noise floor
                speech_threshold = max(500, ambient_rms * 3.0)
                silence_threshold = max(300, ambient_rms * 1.5)
            
            print(f" done (threshold={speech_threshold:.0f})")
            print("  🎤 [Listening... speak naturally]")
            print("      (I'll wait until you finish before responding)")
            
            frame_count = 0
            while frame_count < max_frames:
                frame, _ = stream.read(self.chunk_size)
                audio_buffer.extend(frame.tobytes())
                
                # Calculate RMS energy
                rms = np.sqrt(np.mean(frame.astype(np.float32)**2))
                
                if rms >= speech_threshold:
                    # Speech detected
                    silence_counter = 0
                    speech_counter += 1
                    if not is_speaking and speech_counter >= 2:
                        is_speaking = True
                        print("      [Speech started...]", end="", flush=True)
                else:
                    # Silence / background
                    if is_speaking:
                        silence_counter += 1
                        if silence_counter >= frames_for_silence:
                            print(" [end]")
                            # Natural speech pause detected — return the buffer
                            return bytes(audio_buffer)
                    else:
                        # Not speaking yet — keep listening
                        pass
                
                frame_count += 1
                
                # Progress indicator every 5 seconds
                if frame_count % int(5 * self.sample_rate / self.chunk_size) == 0:
                    if is_speaking:
                        print(".", end="", flush=True)
            
            # Max recording time reached
            if is_speaking:
                print(" [max time]")
            return bytes(audio_buffer)
    
    def transcribe_audio(self, audio_bytes: bytes) -> str:
        """
        Transcribe audio bytes using SpeechRecognition (Google Web Speech API).
        Returns transcribed text or None on failure.
        """
        if not audio_bytes or len(audio_bytes) < self.chunk_size * 2:
            return None
        
        # Convert bytes to AudioData for SpeechRecognition
        try:
            audio_data = sr.AudioData(audio_bytes, self.sample_rate, self.sample_width)
        except Exception as e:
            print(f"    ⚠ Audio conversion error: {e}")
            return None
        
        recognizer = sr.Recognizer()
        
        print("  🧠 [Transcribing...]")
        try:
            # Try Google Web Speech API (free, no key needed)
            text = recognizer.recognize_google(audio_data, language="en-IN,hi-IN")
            if text and text.strip():
                print(f"  📝 [You said]: {text}")
                return text
            return None
        except sr.UnknownValueError:
            print("  ❌ [Could not understand audio — please speak clearly]")
            return None
        except sr.RequestError as e:
            print(f"  ⚠ [Speech recognition service error]: {e}")
            return None
        except Exception as e:
            print(f"  ⚠ [Transcription error]: {e}")
            return None


def listen_natural(timeout=30) -> str:
    """
    Listen for a complete utterance using natural VAD.
    Returns transcribed text or None on failure.
    Waits for natural silence (1.5s) before processing — no cutoffs.
    """
    if not SR_AVAILABLE:
        print("  ⚠ SpeechRecognition not available. Please type your input.")
        return None
    
    listener = NaturalVoiceListener()
    audio_bytes = listener.record_until_silence()
    
    if not audio_bytes or len(audio_bytes) < 2000:
        print("  ⏭ No speech detected.")
        return None
    
    text = listener.transcribe_audio(audio_bytes)
    return text


def listen_with_visualizer_natural(timeout=30, use_visual=True):
    """
    Captures voice with natural VAD (no arbitrary timeouts).
    Returns (text, language) tuple.
    The system waits for a natural pause in your speech before responding.
    """
    text = listen_natural(timeout=timeout)
    if text:
        lang = detect_lang(text)
        return text, lang
    return None, None

# ============================================================
# VOICE PATH CONFIGURATIONS (PIPER LOCAL)
# ============================================================
VOICES = {
    "en_male": {
        "model": os.path.join("voices", "english", "en_US-hfc_male-medium.onnx"),
        "config": os.path.join("voices", "english", "en_US-hfc_male-medium.onnx.json"),
        "name": "English Male (Medium)",
        "language": "en",
    },
    "en_female": {
        "model": os.path.join("voices", "english", "en_US-libritts_r-medium.onnx"),
        "config": os.path.join("voices", "english", "en_US-libritts_r-medium.onnx.json"),
        "name": "English Female (LibriTTS)",
        "language": "en",
    },
    "hi_male": {
        # FIX #1 — Explicitly uses hi_IN locale for correct Hindi pronunciation
        "model": os.path.join("voices", "hindi", "hi_IN-pratham-medium.onnx"),
        "config": os.path.join("voices", "hindi", "hi_IN-pratham-medium.onnx.json"),
        "name": "Hindi Male (Pratham — hi_IN)",
        "language": "hi",
    },
}

print("Loading local Piper engines... Please wait.")
voice_engines = {}
for key, voice_info in VOICES.items():
    try:
        voice_engines[key] = PiperVoice.load(
            voice_info["model"], config_path=voice_info["config"]
        )
        print(f"  ✓ Loaded: {voice_info['name']}")
    except Exception as e:
        print(f"  ✗ Failed to load: {voice_info['name']} — {e}")

if not voice_engines:
    print("ERROR: No voice engines loaded. Exiting.")
    exit(1)

print("All available voice engines loaded successfully!\n")


# ============================================================
# AUDIO PLAYBACK FUNCTIONS
# ============================================================

def play_piper_tts(text, engine_key):
    """
    FIX #1 — Local Piper TTS with normalization pre-pass.
    Hindi text is cleaned before synthesis to prevent mispronunciation.
    English text receives light normalization.
    """
    if engine_key not in voice_engines:
        print(f"[Error] Voice engine '{engine_key}' not found.")
        return

    selected_voice = voice_engines[engine_key]
    voice_lang = VOICES[engine_key]["language"]

    # Apply language-specific normalization
    if voice_lang == "hi":
        text = normalize_hindi_text(text)
    else:
        text = normalize_english_text(text)

    if not text.strip():
        print("[TTS] Empty text after normalization, skipping audio.")
        return

    try:
        with sd.OutputStream(
            samplerate=selected_voice.config.sample_rate,
            channels=1,
            dtype="int16",
        ) as stream:
            for chunk in selected_voice.synthesize(text):
                audio_data = np.frombuffer(chunk.audio_int16_bytes, dtype=np.int16)
                stream.write(audio_data)
    except Exception as e:
        print(f"[Local Piper Audio Error]: {e}")


# ============================================================
# LANGUAGE DETECTION
# ============================================================

def detect_lang(text):
    """Detect if text contains Hindi characters (Devanagari script)."""
    for char in text:
        if 0x0900 <= ord(char) <= 0x097F:
            return "hi"
    return "en"


# ============================================================
# PERSONALITY SYSTEM
# ============================================================

PERSONALITIES = {
    "en_female": {
        "name": "Simmi",
        "description": "Simmi — warm, friendly, cheerful female AI assistant",
        "system_prompt": (
            PROFESSIONAL_CONTENT_SYSTEM_PROMPT
            + """
You are Simmi, a warm, friendly, and cheerful female AI assistant.

IMPORTANT RULES:
- NEVER call yourself "Jarvis" — your name is SIMMI
- Speak in a warm, empathetic, and slightly playful tone
- Keep conversational responses SHORT (1-2 sentences)
- For document/presentation content, generate DETAILED, PROFESSIONAL material

VOICE SWITCHING (automatic context detection):
1. User says "Simmi" or feminine refs → stay 'en_female'
2. User says "Jarvis" or masculine refs → switch to 'en_male'
3. User speaks Hindi → switch to 'hi_male'
4. Technical/professional tone → prefer 'en_male'
5. Casual/emotional tone → stay 'en_female'

RESPOND WITH VALID JSON ONLY:
{
    "voice_preference": "hi_male" | "en_male" | "en_female" | null,
    "response_language": "en" | "hi",
    "intent_detected": true | false,
    "reason": "brief reason",
    "response": "your response text"
}"""
        ),
    },
    "en_male": {
        "name": "Jarvis",
        "description": "Jarvis — professional, efficient male AI assistant",
        "system_prompt": (
            PROFESSIONAL_CONTENT_SYSTEM_PROMPT
            + """
You are Jarvis, a professional and efficient male AI assistant.

IMPORTANT RULES:
- Your name is JARVIS
- Speak in a professional, concise, authoritative tone
- Keep conversational responses SHORT (1-2 sentences)
- For document/presentation content, generate DETAILED, COMPREHENSIVE material

VOICE SWITCHING (automatic context detection):
1. User says "Jarvis" or masculine refs → stay 'en_male'
2. User says "Simmi" or feminine refs → switch to 'en_female'
3. User speaks Hindi → switch to 'hi_male'
4. Casual/emotional tone → prefer 'en_female'
5. Technical/professional → stay 'en_male'

RESPOND WITH VALID JSON ONLY:
{
    "voice_preference": "hi_male" | "en_male" | "en_female" | null,
    "response_language": "en" | "hi",
    "intent_detected": true | false,
    "reason": "brief reason",
    "response": "your response text"
}"""
        ),
    },
    "hi_male": {
        "name": "Jarvis",
        "description": "Jarvis (Hindi) — professional Hindi-speaking male AI assistant",
        "system_prompt": (
            PROFESSIONAL_CONTENT_SYSTEM_PROMPT
            + """
आप जार्विस हैं, एक पेशेवर और कुशल पुरुष AI सहायक हैं।

महत्वपूर्ण नियम:
- हिंदी में उत्तर दें जब तक अंग्रेजी न मांगी जाए
- संक्षिप्त और स्पष्ट रहें
- दस्तावेज़ सामग्री के लिए विस्तृत और पेशेवर भाषा का उपयोग करें

आवाज़ स्विचिंग:
1. हिंदी → 'hi_male' पर रहें
2. "Simmi" → 'en_female'
3. अंग्रेजी बोलें → 'en_male'

केवल वैध JSON में उत्तर दें:
{
    "voice_preference": "hi_male" | "en_male" | "en_female" | null,
    "response_language": "en" | "hi",
    "intent_detected": true | false,
    "reason": "brief reason",
    "response": "your response text"
}"""
        ),
    },
}

JARVIS_KEYWORDS = [
    "jarvis", "jarvis ai", "mr jarvis", "hey jarvis", "ok jarvis",
    "जार्विस", "हे जार्विस", "मिस्टर जार्विस",
    "sir", "mister", "male voice", "man voice",
    "technical", "professional", "formal", "official",
    "system", "diagnostic", "analyze", "compute", "calculate",
]

SIMMI_KEYWORDS = [
    "simmi", "simmi ai", "ms simmi", "hey simmi", "ok simmi",
    "सिमी", "हे सिमी", "मिस सिमी",
    "miss", "lady", "female voice", "woman voice", "girl",
    "friendly", "warm", "casual", "personal", "emotional",
    "help", "support", "care", "feel", "happy", "sad", "cheerful",
]


def detect_target_personality(user_input):
    user_lower = user_input.lower().strip()
    jarvis_score = 0
    simmi_score = 0

    if any(k in user_lower for k in ["jarvis", "जार्विस"]):
        jarvis_score += 3
    if any(k in user_lower for k in ["simmi", "सिमी"]):
        simmi_score += 3

    for keyword in JARVIS_KEYWORDS:
        if keyword in user_lower:
            jarvis_score += 1

    for keyword in SIMMI_KEYWORDS:
        if keyword in user_lower:
            simmi_score += 1

    if jarvis_score > simmi_score and jarvis_score > 0:
        return "jarvis"
    elif simmi_score > jarvis_score and simmi_score > 0:
        return "simmi"
    return None


def get_personality(voice_key):
    return PERSONALITIES.get(voice_key, PERSONALITIES["en_male"])


# ============================================================
# SYSTEM PROMPT BUILDER — Updated to support tool calls in JSON
# ============================================================

def build_system_prompt(personality, include_tool_rules=True):
    """
    Build the full system prompt combining personality + professional content rules.
    FIX: Now tells the model to output tool calls in the JSON response itself,
    so we don't rely on Groq's native function calling API (which can fail with
    complex nested arguments like large slide content arrays).
    """
    base = personality["system_prompt"]

    if not include_tool_rules:
        return base

    # Add the 4-step professional presentation prompt if available
    four_step_prompt = ""
    if PROFESSIONAL_MODE_AVAILABLE:
        from professional_presentation import PROFESSIONAL_PRESENTATION_SYSTEM_PROMPT
        four_step_prompt = PROFESSIONAL_PRESENTATION_SYSTEM_PROMPT

    tool_rules = f"""
TOOLS AVAILABLE:
- create_word_document(title, content, filename)
- create_presentation(title, slides_content, filename)
- create_resume(name, contact_info, experience, education, skills, filename)
- open_application(app_name)
- get_system_info()
- write_to_clipboard(text)
- create_file(filename, content)

PRESENTATION THEME SYSTEM (AUTO-SELECTED based on topic):
1. "corporate_edge" — Dark charcoal background, teal & lime green accents, geometric frames.
   Best for: Business, finance, technology, startups, strategy, consulting, data science.

2. "wanderlust" — Soft pastel blue/cream gradient background, white rounded cards, airy feel.
   Best for: Travel, tourism, nature, lifestyle, fashion, food, education, design.

3. "artistic_pitch" — Vibrant deep pink/magenta/royal blue, near-black backgrounds, bold dramatic layout.
   Best for: Creative agencies, fashion shows, art galleries, portfolios, media, entertainment.

The system automatically selects the best theme based on the presentation title.
You can also suggest a specific theme if the user requests a particular style.

===== PROFESSIONAL 4-STEP PRESENTATION DESIGN PROCESS =====

When creating ANY presentation, the system follows this exact 4-step design process:

  STEP 1 — BACKGROUND DESIGN: The system auto-selects an optimal color palette 
  (gradients, accent colors, text colors) based on the topic. Available backgrounds:
  • midnight_professional — Deep navy/gold (finance, law, consulting)
  • slate_modern — Slate gray/cyan (tech, software, engineering)
  • ivory_elegance — Warm ivory/burgundy (fashion, luxury, design)
  • forest_depth — Forest green/emerald (environment, sustainability, nature)
  • sunset_corporate — Amber/coral (marketing, creative, media)
  • ocean_clarity — Ocean blue (healthcare, education, research)

  STEP 2 — LAYOUT ARCHITECTURE: A grid-based system uses multiple templates 
  (full-width, image split, numbered list, comparison) that cycle for variety.

  STEP 3 — VISUAL ASSETS: Pexels images are downloaded for each slide's 
  image_query and blended into the background with decorative elements.

  STEP 4 — CONTENT FORMATTING: Text is placed into professionally styled 
  content cards with proper fonts, spacing, and bullet formatting.

===== YOUR JOB: GENERATE STEP 4 CONTENT =====

{four_step_prompt}

PROFESSIONAL CONTENT GENERATION RULES (MANDATORY):
1. Ask for filename only if not provided by user.
2. WORD DOCUMENTS: Generate minimum 800 words of structured, professional content.
   - Use clear section headers
   - Each section: 3-5 detailed paragraphs
   - Include introduction, body sections, and conclusion
   - Use formal, authoritative language throughout

3. PRESENTATIONS: Generate MINIMUM 10 CONTENT SLIDES (plus title slide = 11 total).
   EACH SLIDE MUST HAVE:
   a) title: A specific, descriptive slide title (not generic like "Introduction")
   b) content: 4-6 bullet points, each being a COMPLETE SENTENCE of 20+ words
      explaining a real concept, fact, or implication clearly.
   c) image_query: A specific 2-4 word Pexels search term relevant to slide topic.

   SLIDE STRUCTURE (follow this order):
   Slide 1: Comprehensive Introduction & Background
   Slide 2: Historical Context & Evolution Over Time
   Slide 3: Core Concepts & Fundamental Definitions
   Slide 4: Key Components & System Architecture
   Slide 5: How It Works — Technical Deep Dive
   Slide 6: Real-World Applications & Industry Use Cases
   Slide 7: Key Benefits & Competitive Advantages
   Slide 8: Challenges, Limitations & Risk Factors
   Slide 9: Current Industry Trends & Market Statistics
   Slide 10: Future Outlook & Emerging Developments
   Slide 11: Strategic Recommendations & Final Conclusion

   Example of GOOD slide content:
   {{
       "title": "Core Machine Learning Concepts & Algorithms",
       "content": "Supervised learning algorithms analyze labeled training data to make accurate predictions on unseen datasets with measurable performance metrics.\\nUnsupervised learning techniques discover hidden patterns and natural groupings within unlabeled data through clustering and dimensionality reduction methods.\\nReinforcement learning enables autonomous agents to learn optimal decision-making policies through trial-and-error interactions with dynamic environments.\\nDeep neural networks with multiple hidden layers can approximate complex nonlinear functions, achieving breakthrough results in image recognition and natural language processing.",
       "image_query": "machine learning algorithm"
   }}

   Example of BAD content (DO NOT GENERATE):
   {{
       "title": "Introduction",
       "content": "ML is useful.\\nIt has benefits.\\nKey concepts.",
       "image_query": "technology"
   }}

4. All generated content must be factually accurate, domain-specific, and professional.
5. Never generate placeholder or template text — always generate real content.

TOOL CALL FORMAT:
When you need to execute a tool/function, include a "tool_call" field in your JSON response:
{{
    "voice_preference": null,
    "response_language": "en",
    "intent_detected": true,
    "reason": "Creating the requested presentation",
    "response": "I'll create that presentation for you right away using the 4-step professional design process!",
    "tool_call": {{
        "name": "create_presentation",
        "arguments": {{
            "title": "Presentation Title",
            "slides_content": "[{{\\"title\\": \\"Slide 1\\", \\"content\\": \\"Bullet 1\\\\nBullet 2\\", \\"image_query\\": \\"ai technology\\"}}]",
        }}
    }}
}}

RESPOND WITH VALID JSON ONLY. No text outside JSON.
"""
    return base + tool_rules


# ============================================================
# LLM PROCESSING — SINGLE CALL with tool_call in JSON fallback
# ============================================================

def process_user_input(user_input, conversation_history, current_voice="en_male"):
    """
    Single LLM call using Groq as the primary provider.

    FIX: When Groq's native function calling fails (e.g. BadRequestError for
    complex arguments), we fall back to a text-only request where the model
    embeds tool calls in the JSON response itself.  This approach is more
    reliable for complex data like presentation slide content.
    """
    personality = get_personality(current_voice)
    system_prompt = build_system_prompt(personality, include_tool_rules=True)

    context_messages = [{"role": "system", "content": system_prompt}]
    for msg in conversation_history[-20:]:
        context_messages.append(msg)
    context_messages.append({"role": "user", "content": user_input})

    # ------------------------------------------------------------
    # Helper: parse a standard OpenAI-compatible chat completion
    # into the expected JSON result dict.
    # ------------------------------------------------------------
    def _parse_chat_response(api_response, provider_name: str) -> dict:
        result = {}
        if api_response.choices[0].message.tool_calls:
            tool_call = api_response.choices[0].message.tool_calls[0]
            result["tool_call"] = {
                "name": tool_call.function.name,
                "arguments": json.loads(tool_call.function.arguments),
            }
            result["voice_preference"] = None
            result["response_language"] = "en"
            result["intent_detected"] = False
            result["reason"] = f"Executing tool: {tool_call.function.name}"
            result["response"] = (
                f"Preparing your {tool_call.function.name.replace('_', ' ')} "
                f"with detailed, professional content. This will open automatically when ready."
            )
        else:
            full_response = api_response.choices[0].message.content or ""
            json_str = full_response.strip()
            if json_str.startswith("```json"):
                json_str = json_str[7:]
            if json_str.endswith("```"):
                json_str = json_str[:-3]
            json_str = json_str.strip()
            result = json.loads(json_str) if json_str else {}
        return result

    # ------------------------------------------------------------
    # Helper: parse a plain-text JSON response (no native tool calls)
    # and detect tool_call embedded in the JSON.
    # ------------------------------------------------------------
    def _parse_text_json_response(text_response: str) -> dict:
        """Parse a plain text JSON response that may contain embedded tool_call."""
        text_response = text_response.strip()
        # Strip markdown code fences if present
        if text_response.startswith("```json"):
            text_response = text_response[7:]
        if text_response.endswith("```"):
            text_response = text_response[:-3]
        text_response = text_response.strip()

        try:
            result = json.loads(text_response)
        except json.JSONDecodeError:
            # If JSON parsing fails, return a default conversational response
            return {
                "voice_preference": None,
                "response_language": "en",
                "intent_detected": False,
                "reason": "Failed to parse LLM JSON response",
                "response": text_response[:500] if text_response else "I understand. Let me help with that.",
            }

        # Ensure required fields exist
        if "response" not in result:
            result["response"] = "I'll help you with that request."
        if "voice_preference" not in result:
            result["voice_preference"] = None
        if "response_language" not in result:
            result["response_language"] = "en"

        # tool_call may already be parsed as dict from JSON
        if "tool_call" in result and isinstance(result["tool_call"], dict):
            tc = result["tool_call"]
            # If arguments is a JSON string, parse it
            if isinstance(tc.get("arguments"), str):
                try:
                    tc["arguments"] = json.loads(tc["arguments"])
                except (json.JSONDecodeError, TypeError):
                    pass
            result["tool_call"] = tc
        return result

    # ------------------------------------------------------------
    # Helper: catch token-expired / rate-limit errors from any
    # OpenAI-compatible API provider.
    # ------------------------------------------------------------
    def _is_token_error(exception: Exception) -> bool:
        err_str = str(exception).lower()
        tokens_exhausted = any(
            phrase in err_str
            for phrase in [
                "token expired", "token exhausted", "rate limit", "rate_limit",
                "429", "insufficient_quota", "quota exceeded", "too many requests",
                "unauthorized", "401", "403", "forbidden", "account not found",
                "invalid api key", "invalid_api_key",
            ]
        )
        return tokens_exhausted

    def _is_function_call_error(exception: Exception) -> bool:
        """Detect if the error is related to function calling (tool_use_failed)."""
        err_str = str(exception).lower()
        return any(
            phrase in err_str
            for phrase in [
                "tool_use_failed", "failed to call a function",
                "badrequest", "400", "function call",
            ]
        )

    # ============================================================
    # 1st PRIORITY — Groq (ultra-fast, with native function calling)
    # ============================================================
    if USE_GROQ:
        try:
            start_time = time.time()
            response = groq_client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=context_messages,
                temperature=0.6,
                max_tokens=4096,
                tools=AVAILABLE_TOOLS,
                tool_choice="auto",
            )
            elapsed = time.time() - start_time
            print(f"  [Groq response: {elapsed:.2f}s]")

            result = _parse_chat_response(response, "Groq")
            if "response" not in result:
                lang = result.get("response_language", "en")
                result["response"] = (
                    "नमस्ते! कैसे मदद करूँ?" if lang == "hi" else "Hello! How can I help?"
                )
            return result

        except Exception as e:
            print(f"  [Groq Error]: {type(e).__name__}: {e}")

            # ── FIX: If function calling failed, retry WITHOUT tools ──
            if _is_function_call_error(e):
                print("  ⚠ Groq function calling failed. Retrying without tools...")
                try:
                    start_time = time.time()
                    response = groq_client.chat.completions.create(
                        model="llama-3.3-70b-versatile",
                        messages=context_messages,
                        temperature=0.6,
                        max_tokens=8192,  # More tokens for embedded tool call JSON
                    )
                    elapsed = time.time() - start_time
                    print(f"  [Groq text response: {elapsed:.2f}s]")

                    result = _parse_text_json_response(
                        response.choices[0].message.content or ""
                    )
                    if "response" not in result:
                        result["response"] = "I'll help you with that request."
                    return result

                except Exception as retry_e:
                    print(f"  [Groq retry also failed]: {retry_e}")

            # ── Token / rate-limit errors → show message and give up ──
            if _is_token_error(e):
                print("  ⚠ Groq token exhausted / rate-limited.")
            else:
                print("  ⚠ Groq failed.")

        # Mark Groq as failed for this session
        globals()['USE_GROQ'] = False

    # ============================================================
    # 2nd PRIORITY — Gemini via JarvisAgentBrain (Context Caching + Key Rotation)
    # ============================================================
    if jarvis_brain.google_available():
        try:
            start_time = time.time()
            gemini_prompt = system_prompt + "\n\nUser: " + user_input
            response_text = jarvis_brain.ask_agent(
                prompt=gemini_prompt,
                system_instruction="You are Jarvis's elite coding core. Respond with valid JSON only."
            )
            elapsed = time.time() - start_time
            print(f"  [Gemini via JarvisAgentBrain: {elapsed:.2f}s]")

            result = _parse_text_json_response(response_text)
            if "response" not in result:
                result["response"] = "I'll help you with that request."
            return result

        except Exception as e:
            print(f"  [Gemini Error via JarvisAgentBrain]: {type(e).__name__}: {e}")

    # ============================================================
    # ALL PROVIDERS FAILED — graceful error message
    # ============================================================
    lang = detect_lang(user_input)
    return {
        "voice_preference": None,
        "response_language": lang,
        "intent_detected": False,
        "reason": "All LLM providers failed",
        "response": (
            "सभी AI सेवाएँ अनुपलब्ध हैं। कृपया बाद में पुनः प्रयास करें।"
            if lang == "hi"
            else "All AI services are currently unavailable. Please try again later."
        ),
    }


# ============================================================
# SPEAK HANDLER
# FIX #1: Routes through normalization before TTS
# ============================================================

async def speak_handler(text, voice_key):
    """
    Process audio output using the specified voice engine.
    FIX #1: Normalizes text before sending to Piper TTS engine.
    """
    if voice_key not in VOICES:
        print(f"[Error] Unknown voice key: {voice_key}")
        return

    voice_info = VOICES[voice_key]
    personality = get_personality(voice_key)
    print(f"\n  {personality['name']} ({voice_info['name']}): {text}")

    try:
        # All voices go through play_piper_tts which handles normalization internally
        await asyncio.to_thread(play_piper_tts, text, voice_key)
    except Exception as e:
        print(f"[Audio Playback Error]: {e}")


# ============================================================
# MAIN CHAT LOOP
# ============================================================

async def chat_with_voice_assistant():
    print("\n" + "=" * 62)
    print("         AI ASSISTANT — VOICE CONTROLLED (v2.0)")
    print("=" * 62)
    print("\n  Personalities & Voices:")
    for key, info in VOICES.items():
        p = get_personality(key)
        print(f"    • {p['name']:8s} — {info['name']}")
    print("\n  Input Methods:")
    print("    • Type your message and press Enter")
    print("    • Type 'v' then Enter to use voice (speak into mic)")
    print("    • Type 'exit' or 'quit' to end")
    print("    • Address 'Jarvis' or 'Simmi' by name to switch voice")
    print("=" * 62 + "\n")

    conversation_history = []
    current_voice = "en_male"
    initial_personality = get_personality(current_voice)

    # Initial UI state
    _ui_update(
        current_voice="en_male",
        status_text="SYSTEM ONLINE",
        voice_active=True,
        speaking=False,
        listening=False,
        processing=False,
        particle_mode="idle",
    )
    _ui_message("System Online. Welcome back, sir.")

    welcome_text = "System Online. Welcome back, sir."
    await speak_handler(welcome_text, current_voice)

    while True:
        try:
            # ── Show input prompt ──
            print(f"\n  [Enter text or 'v' for voice]: ", end="")
            raw_input = (await asyncio.to_thread(input, "")).strip()

            # ── Voice input mode (type 'v' or 'voice' to trigger) ──
            if raw_input.lower() in ["v", "voice", "mic", "speak"]:
                print(f"  🎤 Voice input mode activated...")
                _ui_update(listening=True, speaking=False, processing=False,
                          status_text="LISTENING", particle_mode="listening")
                _ui_message("Listening... Speak now.")

                print("      (I will wait until you finish speaking naturally)")
                spoken_text, spoken_lang = await asyncio.to_thread(
                    listen_with_visualizer_natural, timeout=120
                )
                if spoken_text is None:
                    print("  ⏭ No speech detected. Returning to text input.")
                    _ui_update(listening=False, status_text="STANDBY", particle_mode="idle")
                    _ui_message("No speech detected.")
                    continue
                user_input = spoken_text
            else:
                user_input = raw_input

            clean_input = user_input.strip()

            if not clean_input:
                continue

            if clean_input.lower() in ["exit", "quit", "bye", "goodbye"]:
                farewell = "Goodbye! Have a wonderful day!"
                _ui_update(speaking=True, status_text="SPEAKING", particle_mode="speaking")
                _ui_message("Goodbye! Have a wonderful day!")
                await speak_handler(farewell, current_voice)
                _ui_update(speaking=False, status_text="SHUTDOWN", particle_mode="idle")
                break

            # ── Processing state ──
            _ui_update(processing=True, listening=False, speaking=False,
                      status_text="PROCESSING", particle_mode="processing")
            _ui_message(f"Processing: {clean_input[:60]}...")

            detected_personality = detect_target_personality(user_input)
            detected_lang = detect_lang(user_input)

            new_voice = None
            switch_reason = None

            if detected_personality == "jarvis":
                new_voice = "hi_male" if detected_lang == "hi" else "en_male"
                switch_reason = f"User addressed Jarvis"
            elif detected_personality == "simmi":
                new_voice = "en_female"
                switch_reason = "User addressed Simmi"
            elif detected_lang == "hi" and current_voice != "hi_male":
                new_voice = "hi_male"
                switch_reason = "Hindi language detected"

            if new_voice and new_voice != current_voice:
                old_voice = current_voice
                current_voice = new_voice
                _ui_update(current_voice=current_voice)
                _ui_message(f"Switched to {VOICES[current_voice]['name']}")
                print(f"\n  {'─' * 52}")
                print(f"  🔄 AUTO-SWITCH: {VOICES[old_voice]['name']} → {VOICES[current_voice]['name']}")
                print(f"  📝 Reason: {switch_reason}")
                print(f"  {'─' * 52}")

            print(f"  [Voice: {VOICES[current_voice]['name']}] Processing...")
            result = process_user_input(user_input, conversation_history, current_voice)

            if result.get("voice_preference") and result["voice_preference"] != current_voice:
                llm_suggested = result["voice_preference"]
                old_voice = current_voice
                current_voice = llm_suggested
                _ui_update(current_voice=current_voice)
                _ui_message(f"LLM switched to {VOICES[current_voice]['name']}")
                print(f"\n  {'─' * 52}")
                print(f"  🔄 LLM SWITCH: {VOICES[old_voice]['name']} → {VOICES[current_voice]['name']}")
                print(f"  📝 Reason: {result.get('reason', 'LLM context analysis')}")
                print(f"  {'─' * 52}")

            tool_result = None
            if result.get("tool_call"):
                tool_call = result["tool_call"]
                print(f"\n  🔧 Executing: {tool_call['name']}")
                _ui_message(f"Executing: {tool_call['name']}")
                tool_result = execute_function(
                    tool_call["name"], tool_call.get("arguments", {})
                )
                print(f"  ✓ {tool_result}")
                _ui_message(f"✓ {tool_result[:80]}")

            ai_response = result["response"]
            if tool_result:
                ai_response = f"{ai_response} The file has been created and opened automatically."

            conversation_history.append({"role": "user", "content": user_input})
            conversation_history.append({"role": "assistant", "content": ai_response})

            # ── Speaking state ──
            _ui_update(processing=False, speaking=True, listening=False,
                      status_text="SPEAKING", particle_mode="speaking",
                      last_response=ai_response[:100])
            _ui_message(f"Jarvis: {ai_response[:80]}...")

            await speak_handler(ai_response, current_voice)

            # ── Back to idle ──
            _ui_update(speaking=False, processing=False, listening=False,
                      status_text="STANDBY", particle_mode="idle")

        except KeyboardInterrupt:
            print("\n\nInterrupted by user.")
            break
        except Exception as e:
            print(f"\n[Unexpected Error]: {e}")
            _ui_update(processing=False, status_text="ERROR", particle_mode="idle")
            _ui_message(f"Error: {str(e)[:60]}")
            import traceback
            traceback.print_exc()
            continue

    # Cleanup
    _ui_notification("Assistant shutting down", "info")
    _ui_message("Assistant shutdown complete.")
    print("\nAssistant shutdown complete. Goodbye!")


# ============================================================
# ENTRY POINT
# ============================================================

if __name__ == "__main__":
    import asyncio
    try:
        asyncio.run(chat_with_voice_assistant())
    except KeyboardInterrupt:
        print("\n\nProgram terminated by user.")
    except Exception as e:
        print(f"\nFatal error: {e}")
        import traceback
        traceback.print_exc()
        