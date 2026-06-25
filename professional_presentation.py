"""
Professional Presentation Engine (v4.0) — Step-by-Step Design Process
=====================================================================
Implements the exact 4-step presentation design process:

  Step 1: Background Design
  Step 2: Professional Layout Architecture
  Step 3: Visual Assets Integration
  Step 4: Content Generation & Formatting

This engine produces presentation-ready content and slides that follow
a methodical, professional workflow — exactly as specified by the user.
"""

import os
import json
import time
import re
import random
from datetime import datetime
from typing import List, Dict, Optional, Tuple, Any
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

from pexels_integration import (
    get_pexels_client,
    download_background_image,
    get_color_palette_from_image,
    create_background_with_overlay,
    MEDIA_DIR,
)


OUTPUT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "generated_files")
os.makedirs(OUTPUT_DIR, exist_ok=True)


# ══════════════════════════════════════════════════════════════════════════════
# STEP 1: BACKGROUND DESIGN
# ══════════════════════════════════════════════════════════════════════════════
# Defines the visual foundation — color palettes, gradients, textures.

@dataclass
class BackgroundDesign:
    """
    Complete background design specification.
    Covers color palette, gradient direction, textures, and overlay treatment.
    """
    name: str
    description: str
    base_color: Tuple[int, int, int]
    secondary_color: Tuple[int, int, int]
    accent_color_1: Tuple[int, int, int]
    accent_color_2: Tuple[int, int, int]
    text_primary: Tuple[int, int, int]
    text_secondary: Tuple[int, int, int]
    text_accent: Tuple[int, int, int]
    gradient_direction: str = "top_to_bottom"  # or "diagonal", "left_to_right"
    has_texture_overlay: bool = False
    overlay_opacity: int = 35  # 0-100
    use_gradient: bool = True
    card_bg_color: Tuple[int, int, int] = None
    card_border_color: Tuple[int, int, int] = None

    def __post_init__(self):
        if self.card_bg_color is None:
            # Derive a slightly lighter/darker card color from base
            r = min(255, self.base_color[0] + 8)
            g = min(255, self.base_color[1] + 8)
            b = min(255, self.base_color[2] + 8)
            self.card_bg_color = (r, g, b)
        if self.card_border_color is None:
            self.card_border_color = self.accent_color_1


# Pre-built professional color palettes
PROFESSIONAL_BACKGROUNDS = {
    "midnight_professional": BackgroundDesign(
        name="midnight_professional",
        description="Deep navy-to-charcoal gradient with gold accents — authoritative, trustworthy",
        base_color=(18, 22, 38),
        secondary_color=(30, 35, 55),
        accent_color_1=(200, 170, 110),   # Gold
        accent_color_2=(100, 130, 200),   # Steel blue
        text_primary=(255, 255, 255),
        text_secondary=(180, 185, 195),
        text_accent=(200, 170, 110),
        gradient_direction="diagonal",
        card_bg_color=(25, 30, 48),
        card_border_color=(200, 170, 110),
    ),
    "slate_modern": BackgroundDesign(
        name="slate_modern",
        description="Slate gray to dark blue-gray gradient with cyan accents — modern, tech-forward",
        base_color=(28, 32, 40),
        secondary_color=(38, 45, 58),
        accent_color_1=(0, 200, 220),     # Cyan
        accent_color_2=(80, 160, 200),    # Light blue
        text_primary=(255, 255, 255),
        text_secondary=(190, 200, 210),
        text_accent=(0, 200, 220),
        gradient_direction="top_to_bottom",
        card_bg_color=(34, 40, 50),
        card_border_color=(0, 200, 220),
    ),
    "ivory_elegance": BackgroundDesign(
        name="ivory_elegance",
        description="Warm ivory to cream gradient with burgundy accents — elegant, sophisticated",
        base_color=(245, 240, 230),
        secondary_color=(235, 225, 210),
        accent_color_1=(140, 50, 60),     # Burgundy
        accent_color_2=(190, 150, 110),   # Bronze
        text_primary=(35, 30, 28),
        text_secondary=(100, 95, 90),
        text_accent=(140, 50, 60),
        gradient_direction="left_to_right",
        card_bg_color=(255, 252, 248),
        card_border_color=(140, 50, 60),
    ),
    "forest_depth": BackgroundDesign(
        name="forest_depth",
        description="Deep forest green to dark teal gradient with emerald accents — natural, growth-focused",
        base_color=(10, 30, 20),
        secondary_color=(15, 45, 35),
        accent_color_1=(60, 180, 110),    # Emerald
        accent_color_2=(80, 160, 130),    # Sage
        text_primary=(240, 245, 240),
        text_secondary=(170, 190, 175),
        text_accent=(60, 180, 110),
        gradient_direction="top_to_bottom",
        card_bg_color=(16, 38, 28),
        card_border_color=(60, 180, 110),
    ),
    "sunset_corporate": BackgroundDesign(
        name="sunset_corporate",
        description="Warm amber-to-charcoal gradient with coral accents — energetic, creative business",
        base_color=(30, 25, 30),
        secondary_color=(50, 35, 40),
        accent_color_1=(230, 120, 80),    # Coral
        accent_color_2=(210, 160, 80),    # Amber
        text_primary=(255, 250, 245),
        text_secondary=(200, 190, 180),
        text_accent=(230, 120, 80),
        gradient_direction="diagonal",
        card_bg_color=(40, 33, 38),
        card_border_color=(230, 120, 80),
    ),
    "ocean_clarity": BackgroundDesign(
        name="ocean_clarity",
        description="Light ocean blue to white gradient with deep blue accents — clean, clear, trustworthy",
        base_color=(220, 235, 250),
        secondary_color=(240, 245, 255),
        accent_color_1=(30, 80, 160),     # Deep blue
        accent_color_2=(70, 150, 200),    # Sky blue
        text_primary=(20, 30, 50),
        text_secondary=(80, 95, 120),
        text_accent=(30, 80, 160),
        gradient_direction="top_to_bottom",
        card_bg_color=(255, 255, 255),
        card_border_color=(30, 80, 160),
    ),
}


def auto_select_background(topic: str) -> BackgroundDesign:
    """
    Automatically select the best background design based on topic keywords.
    This is STEP 1 of the presentation design process.
    """
    topic_lower = topic.lower()
    
    scoring = {
        "midnight_professional": 0,
        "slate_modern": 0,
        "ivory_elegance": 0,
        "forest_depth": 0,
        "sunset_corporate": 0,
        "ocean_clarity": 0,
    }
    
    # Keyword mapping
    keywords = {
        "midnight_professional": [
            "finance", "banking", "investment", "law", "legal", "consulting",
            "executive", "board", "governance", "audit", "risk", "compliance",
            "corporate strategy", "merger", "acquisition", "wealth",
        ],
        "slate_modern": [
            "technology", "software", "saas", "startup", "engineering",
            "data science", "ai", "machine learning", "cloud", "cybersecurity",
            "devops", "agile", "digital transformation", "innovation",
            "programming", "architecture", "system design",
        ],
        "ivory_elegance": [
            "fashion", "luxury", "design", "art", "gallery", "exhibition",
            "wedding", "event", "hospitality", "fine dining", "wine",
            "premium", "exclusive", "boutique", "heritage", "craftsmanship",
        ],
        "forest_depth": [
            "environment", "sustainability", "ecology", "climate", "nature",
            "conservation", "renewable", "green", "organic", "agriculture",
            "wildlife", "forest", "ocean", "clean energy", "biodiversity",
        ],
        "sunset_corporate": [
            "marketing", "advertising", "creative agency", "branding",
            "media", "entertainment", "social media", "content",
            "digital marketing", "growth", "campaign", "creative",
        ],
        "ocean_clarity": [
            "healthcare", "medicine", "education", "research", "science",
            "academic", "university", "training", "coaching", "wellness",
            "mental health", "public health", "pharmaceutical", "biotech",
        ],
    }
    
    for bg_name, kw_list in keywords.items():
        for kw in kw_list:
            if kw in topic_lower:
                scoring[bg_name] += 2
            for word in topic_lower.split():
                if word == kw or word.startswith(kw) or kw.startswith(word):
                    scoring[bg_name] += 1
    
    # If no clear match, use generic heuristics
    if max(scoring.values()) == 0:
        return PROFESSIONAL_BACKGROUNDS["slate_modern"]  # Safe modern default
    
    best_name = max(scoring, key=scoring.get)
    return PROFESSIONAL_BACKGROUNDS[best_name]


# ══════════════════════════════════════════════════════════════════════════════
# STEP 2: PROFESSIONAL LAYOUT ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
# Defines the grid system and slide layouts.

@dataclass
class SlideLayoutSpec:
    """
    Defines the exact positioning and grid structure for a slide type.
    All measurements are in inches (will be converted to EMU at render time).
    """
    name: str
    description: str
    
    # Margins
    margin_left: float = 0.6
    margin_right: float = 0.6
    margin_top: float = 0.4
    margin_bottom: float = 0.5
    
    # Header area
    header_height: float = 0.9
    header_left: float = 0.6
    header_width: float = 12.1
    
    # Accent bar
    accent_bar_top: float = 1.3
    accent_bar_height: float = 0.04
    accent_bar_width: float = 2.0
    accent_bar_left: float = 0.6
    
    # Content area
    content_top: float = 1.6
    content_left: float = 0.7
    content_width: float = 11.9
    content_height: float = 5.5
    
    # Image placeholder (if applicable)
    image_left: float = None # Left coordinate for side-by-side images
    image_right: float = 7.5
    image_top: float = 1.6
    image_width: float = 5.2
    image_height: float = 5.5
    
    # Title slide specific
    title_top: float = 2.5
    title_center_x: float = 6.67  # Center of 13.333" slide
    title_width: float = 10.0
    subtitle_top: float = 4.5
    subtitle_width: float = 8.0
    
    # Number badge
    number_badge_left: float = 0.3
    number_badge_top: float = 1.5
    number_badge_size: float = 0.7


# Standard layout templates
LAYOUTS = {
    "title_slide": SlideLayoutSpec(
        name="title_slide",
        description="Center-aligned title with subtitle and decorative elements",
        margin_top=0.0,
        margin_left=0.0,
        title_top=2.5,
        title_width=10.0,
        subtitle_top=4.5,
        subtitle_width=8.0,
    ),
    "content_full": SlideLayoutSpec(
        name="content_full",
        description="Full-width content with header bar, accent line, and text area",
        header_height=0.9,
        content_top=1.6,
        content_width=11.9,
        content_height=5.5,
    ),
    "content_with_image_right": SlideLayoutSpec(
        name="content_with_image_right",
        description="Text on left (60%), image on right (40%)",
        content_left=0.7,
        content_width=6.5, # Reduced from 6.8 to avoid bleed or overlap
        content_height=5.5,
        image_right=7.7, # Starts clearly after text column (0.7 + 6.5 = 7.2 + 0.5 gap = 7.7)
        image_top=1.6,
        image_width=5.0,
        image_height=5.5,
    ),
    "content_with_image_left": SlideLayoutSpec(
        name="content_with_image_left",
        description="Image on left (40%), text on right (60%)",
        content_left=6.3, # Starts after image column (0.7 + 5.0 = 5.7 + 0.6 gap = 6.3)
        content_width=6.4,
        content_height=5.5,
        image_left=0.7, # Positioned explicitly on the left
        image_right=0.7,
        image_top=1.6,
        image_width=5.0,
        image_height=5.5,
    ),
    "data_visualization": SlideLayoutSpec(
        name="data_visualization",
        description="Header with extra body space for charts and data",
        header_height=0.8,
        content_top=1.4,
        content_width=11.9,
        content_height=5.8,
    ),
    "numbered_list": SlideLayoutSpec(
        name="numbered_list",
        description="Content with large number badges on the left",
        header_height=0.9,
        content_top=1.6,
        content_left=1.3,
        content_width=11.3,
        content_height=5.5,
        number_badge_left=0.3,
        number_badge_top=1.5,
        number_badge_size=0.7,
    ),
    "split_comparison": SlideLayoutSpec(
        name="split_comparison",
        description="Two-column comparison layout",
        content_left=0.7,
        content_width=5.5,
        content_height=5.5,
        image_right=7.0,
        image_top=1.6,
        image_width=5.8,
        image_height=5.5,
    ),
}


def select_best_layout(slide_index: int, total_slides: int, has_image: bool = True) -> SlideLayoutSpec:
    """
    Select the best layout for a slide based on its position and content type.
    This ensures visual variety throughout the presentation.
    """
    # Title slide always uses title_slide layout
    if slide_index == 0:
        return LAYOUTS["title_slide"]
    
    # Cycle through layouts for visual variety
    layout_cycle = [
        "content_full",
        "content_with_image_right",
        "content_full",
        "numbered_list",
        "content_with_image_left",
        "content_full",
        "split_comparison",
        "content_with_image_right",
        "content_full",
        "numbered_list",
        "data_visualization",
        "content_with_image_left",
        "content_full",
    ]
    
    # Use position in cycle
    idx = (slide_index - 1) % len(layout_cycle)
    layout_name = layout_cycle[idx]
    return LAYOUTS[layout_name]


# ══════════════════════════════════════════════════════════════════════════════
# OOXML HELPERS — for fills, borders, and effects
# ══════════════════════════════════════════════════════════════════════════════

def _apply_solid_fill(shape, r: int, g: int, b: int):
    """Apply a solid (opaque) fill to a shape via OOXML."""
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        spPr = etree.SubElement(sp, qn("p:spPr"))
    for tag in ["a:noFill", "a:solidFill", "a:gradFill", "a:pattFill", "a:blipFill"]:
        for child in spPr.findall(qn(tag)):
            spPr.remove(child)
    hex_color = f"{r:02X}{g:02X}{b:02X}"
    solid_fill = etree.SubElement(spPr, qn("a:solidFill"))
    srgb_clr = etree.SubElement(solid_fill, qn("a:srgbClr"))
    srgb_clr.set("val", hex_color)


def _apply_soft_gradient_fill(shape, color1: Tuple[int,int,int], color2: Tuple[int,int,int], angle: int = 5400000):
    """Apply a linear gradient fill with configurable angle."""
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        spPr = etree.SubElement(sp, qn("p:spPr"))
    for tag in ["a:noFill", "a:solidFill", "a:gradFill", "a:pattFill", "a:blipFill"]:
        for child in spPr.findall(qn(tag)):
            spPr.remove(child)

    grad = etree.SubElement(spPr, qn("a:gradFill"))
    gs_lst = etree.SubElement(grad, qn("a:gsLst"))

    gs1 = etree.SubElement(gs_lst, qn("a:gs"))
    gs1.set("pos", "0")
    c1 = etree.SubElement(gs1, qn("a:srgbClr"))
    c1.set("val", f"{color1[0]:02X}{color1[1]:02X}{color1[2]:02X}")

    gs2 = etree.SubElement(gs_lst, qn("a:gs"))
    gs2.set("pos", "100000")
    c2 = etree.SubElement(gs2, qn("a:srgbClr"))
    c2.set("val", f"{color2[0]:02X}{color2[1]:02X}{color2[2]:02X}")

    lin = etree.SubElement(grad, qn("a:lin"))
    lin.set("ang", str(angle))
    lin.set("scaled", "1")


def _apply_blended_fill(shape, r: int, g: int, b: int, alpha_percent: int):
    """Apply semi-transparent fill (glassmorphism effect)."""
    sp = shape._element
    hex_color = f"{r:02X}{g:02X}{b:02X}"
    alpha_val = int((alpha_percent / 100) * 100000)
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        spPr = etree.SubElement(sp, qn("p:spPr"))
    for tag in ["a:noFill", "a:solidFill", "a:gradFill", "a:pattFill", "a:blipFill"]:
        for child in spPr.findall(qn(tag)):
            spPr.remove(child)
    solid_fill = etree.SubElement(spPr, qn("a:solidFill"))
    srgb_clr = etree.SubElement(solid_fill, qn("a:srgbClr"))
    srgb_clr.set("val", hex_color)
    alpha_elem = etree.SubElement(srgb_clr, qn("a:alpha"))
    alpha_elem.set("val", str(100000 - alpha_val))


def _remove_shape_border(shape):
    """Remove border from a shape."""
    try:
        shape.line.fill.background()
    except Exception:
        pass


def _set_shape_border(shape, r: int, g: int, b: int, width: float = 1.0):
    """Set a colored border on a shape via OOXML."""
    sp = shape._element
    hex_color = f"{r:02X}{g:02X}{b:02X}"
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        spPr = etree.SubElement(sp, qn("p:spPr"))
    for child in spPr.findall(qn("a:ln")):
        spPr.remove(child)
    ln = etree.SubElement(spPr, qn("a:ln"))
    ln.set("w", str(int(width * 12700)))
    sf = etree.SubElement(ln, qn("a:solidFill"))
    sc = etree.SubElement(sf, qn("a:srgbClr"))
    sc.set("val", hex_color)


def _set_shape_shadow(shape, r: int, g: int, b: int, blur_radius: int = 80000):
    """Add a soft shadow to a shape."""
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        spPr = etree.SubElement(sp, qn("p:spPr"))
    for child in spPr.findall(qn("a:effectLst")):
        spPr.remove(child)
    effect_lst = etree.SubElement(spPr, qn("a:effectLst"))
    outer_shdw = etree.SubElement(effect_lst, qn("a:outerShdw"))
    outer_shdw.set("blurRad", str(blur_radius))
    outer_shdw.set("dist", "40000")
    outer_shdw.set("dir", "2700000")
    outer_shdw.set("algn", "tl")
    outer_shdw.set("rotWithShape", "0")
    srgb = etree.SubElement(outer_shdw, qn("a:srgbClr"))
    srgb.set("val", f"{r:02X}{g:02X}{b:02X}")
    alpha = etree.SubElement(srgb, qn("a:alpha"))
    alpha.set("val", "40000")


# ══════════════════════════════════════════════════════════════════════════════
# PROFESSIONAL SLIDE DESIGNER — Implements Steps 1-3
# ══════════════════════════════════════════════════════════════════════════════

class ProfessionalSlideDesigner:
    """
    Designs slides following the 4-step professional process.
    
    Step 1: Background Design — applies color palette and gradients
    Step 2: Layout Architecture — positions elements on a grid
    Step 3: Visual Assets Integration — adds images and decorative elements
    Step 4: Content Formatting — formats text within the layout
    """
    
    def __init__(self, background: BackgroundDesign, slide_width: float = 13.333, slide_height: float = 7.5):
        self.bg = background
        self.slide_width = Inches(slide_width)
        self.slide_height = Inches(slide_height)
        self._slide_width_inches = slide_width
        self._slide_height_inches = slide_height
        self.pexels_client = get_pexels_client()
        
        # Step 2: Load default layout
        self.current_layout = LAYOUTS["content_full"]
    
    def rgb_to_pptx(self, rgb: Tuple[int, int, int]) -> RGBColor:
        return RGBColor(rgb[0], rgb[1], rgb[2])
    
    def _inches(self, val: float) -> Inches:
        return Inches(val)
    
    # ──────────────────────────────────────────────────────────────────────────
    # STEP 1: BACKGROUND DESIGN — Apply the visual foundation
    # ──────────────────────────────────────────────────────────────────────────
    
    def apply_background(self, slide) -> None:
        """
        STEP 1: Apply the background design to a slide.
        Uses gradient fills based on the selected background specification.
        """
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            width=self.slide_width, height=self.slide_height,
        )
        bg_shape.name = "BackgroundDesign"
        
        # Determine gradient angle based on direction
        angle_map = {
            "top_to_bottom": 5400000,
            "left_to_right": 0,
            "diagonal": 2700000,
        }
        angle = angle_map.get(self.bg.gradient_direction, 5400000)
        
        _apply_soft_gradient_fill(
            bg_shape,
            self.bg.base_color,
            self.bg.secondary_color,
            angle=angle,
        )
        _remove_shape_border(bg_shape)
        
        # Send to back
        sp = bg_shape.shape._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)
        
        # Add subtle decorative accent bar at the very top
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            width=self.slide_width, height=Pt(3),
        )
        _apply_solid_fill(accent_bar, self.bg.accent_color_1[0], self.bg.accent_color_1[1], self.bg.accent_color_1[2])
        _remove_shape_border(accent_bar)
        sp = accent_bar.shape._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)
    
    # ──────────────────────────────────────────────────────────────────────────
    # STEP 2: LAYOUT ARCHITECTURE — Build the grid structure
    # ──────────────────────────────────────────────────────────────────────────
    
    def set_layout(self, layout_spec: SlideLayoutSpec) -> None:
        """Set the current layout specification (Step 2)."""
        self.current_layout = layout_spec
    
    def add_header_bar(self, slide, title: str) -> None:
        """Add a styled header/title bar to the slide."""
        layout = self.current_layout
        
        # Header background strip
        header_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self._inches(layout.margin_left),
            self._inches(layout.margin_top),
            self._inches(layout.header_width),
            self._inches(layout.header_height),
        )
        _apply_blended_fill(
            header_bg,
            self.bg.accent_color_1[0], self.bg.accent_color_1[1], self.bg.accent_color_1[2],
            88,  # 12% opacity
        )
        _remove_shape_border(header_bg)
        
        # Title text
        tbox = slide.shapes.add_textbox(
            self._inches(layout.margin_left + 0.15),
            self._inches(layout.margin_top + 0.08),
            self._inches(layout.header_width - 0.3),
            self._inches(layout.header_height - 0.16),
        )
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.rgb_to_pptx(self.bg.text_primary)
        p.font.name = "Calibri"
    
    def add_accent_line(self, slide) -> None:
        """Add a decorative accent line below the header."""
        layout = self.current_layout
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self._inches(layout.accent_bar_left),
            self._inches(layout.accent_bar_top),
            self._inches(layout.accent_bar_width),
            self._inches(layout.accent_bar_height),
        )
        _apply_solid_fill(line, self.bg.accent_color_1[0], self.bg.accent_color_1[1], self.bg.accent_color_1[2])
        _remove_shape_border(line)
        
        # Second thinner line in accent color 2
        line2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self._inches(layout.accent_bar_left + layout.accent_bar_width + 0.15),
            self._inches(layout.accent_bar_top + 0.005),
            self._inches(0.8),
            self._inches(layout.accent_bar_height * 0.6),
        )
        _apply_solid_fill(line2, self.bg.accent_color_2[0], self.bg.accent_color_2[1], self.bg.accent_color_2[2])
        _remove_shape_border(line2)
    
    def add_number_badge(self, slide, number: int) -> None:
        """Add a large number badge for numbered slides."""
        layout = self.current_layout
        num = slide.shapes.add_textbox(
            self._inches(layout.number_badge_left),
            self._inches(layout.number_badge_top),
            self._inches(layout.number_badge_size),
            self._inches(layout.number_badge_size),
        )
        tf = num.text_frame
        p = tf.paragraphs[0]
        p.text = str(number).zfill(2)
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(42)
        p.font.bold = True
        p.font.color.rgb = self.rgb_to_pptx(self.bg.accent_color_1)
        p.font.name = "Calibri"
    
    # ──────────────────────────────────────────────────────────────────────────
    # STEP 3: VISUAL ASSETS INTEGRATION — Add images and decorative elements
    # ──────────────────────────────────────────────────────────────────────────
    
    def integrate_background_image(self, slide, image_path: str, overlay_opacity: int = None) -> None:
        """
        STEP 3: Integrate a background image with proper overlay for readability.
        The image is placed as a full-bleed background with a semi-transparent
        color overlay on top.
        """
        if not image_path or not os.path.exists(image_path):
            return
        
        try:
            # Add the image as full-bleed background
            slide.shapes.add_picture(
                image_path, 0, 0,
                width=self.slide_width, height=self.slide_height,
            )
            
            # Add semi-transparent overlay for text readability
            overlay = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0,
                width=self.slide_width, height=self.slide_height,
            )
            op = overlay_opacity if overlay_opacity is not None else self.bg.overlay_opacity
            _apply_blended_fill(
                overlay,
                self.bg.base_color[0], self.bg.base_color[1], self.bg.base_color[2],
                op,
            )
            _remove_shape_border(overlay)
            
            # Send overlay to front of background but behind content
            sp = overlay.shape._element
            sp.getparent().remove(sp)
            slide.shapes._spTree.insert(2, sp)
            
        except Exception as e:
            print(f"     ⚠ Image integration skipped: {e}")
    
    def add_decorative_element(self, slide) -> None:
        """
        Add a subtle decorative geometric element to enhance visual interest.
        Part of Step 3: Visual Assets Integration.
        """
        # Add a subtle circle in the corner
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            self._inches(self._slide_width_inches - 1.5),
            self._inches(self._slide_height_inches - 1.5),
            self._inches(1.2), self._inches(1.2),
        )
        _apply_blended_fill(
            circle,
            self.bg.accent_color_1[0], self.bg.accent_color_1[1], self.bg.accent_color_1[2],
            92,  # 8% opacity
        )
        _remove_shape_border(circle)
        
        # Send to back (below content)
        sp = circle.shape._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)
    
    # ──────────────────────────────────────────────────────────────────────────
    # STEP 4: CONTENT FORMATTING — Format text within the layout
    # ──────────────────────────────────────────────────────────────────────────
    
    def format_content_text(self, slide, title: str, content: str, slide_index: int = 0) -> None:
        """
        STEP 4: Format and place content text within the established layout.
        Creates a styled content card with title and bullet points.
        """
        layout = self.current_layout
        
        # ── Content card (background) ─────────────────────────────────────────
        card_left = self._inches(layout.content_left - 0.1)
        card_top = self._inches(layout.content_top - 0.1)
        card_width = self._inches(layout.content_width + 0.2)
        card_height = self._inches(layout.content_height + 0.2)
        
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            card_left, card_top, card_width, card_height,
        )
        _apply_solid_fill(card, self.bg.card_bg_color[0], self.bg.card_bg_color[1], self.bg.card_bg_color[2])
        _set_shape_border(card, self.bg.card_border_color[0], self.bg.card_border_color[1], self.bg.card_border_color[2], 0.5)
        _set_shape_shadow(card, 0, 0, 0, blur_radius=60000)
        
        # ── Content text box ──────────────────────────────────────────────────
        tbox = slide.shapes.add_textbox(
            self._inches(layout.content_left),
            self._inches(layout.content_top),
            self._inches(layout.content_width),
            self._inches(layout.content_height),
        )
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.15)
        tf.margin_bottom = Inches(0.1)
        
        # Process content into formatted bullet points
        lines = content.split("\n") if content else []
        
        first_line = True
        for line in lines:
            stripped = line.strip()
            if not stripped:
                continue
            
            # Clean leading bullets
            if stripped.startswith(("•", "-", "*", "→", "▸", "◦", "—")):
                stripped = stripped[1:].strip()
            
            # If it looks like a sub-heading (SHORT, no period), format differently
            is_subheading = len(stripped) < 60 and not stripped.endswith(("."))
            
            if first_line:
                p = tf.paragraphs[0]
                first_line = False
            else:
                p = tf.add_paragraph()
            
            if is_subheading and len(stripped) < 50:
                # Sub-heading style
                p.text = stripped
                p.alignment = PP_ALIGN.LEFT
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = self.rgb_to_pptx(self.bg.accent_color_1)
                p.font.name = "Calibri"
                p.space_after = Pt(6)
                p.space_before = Pt(8)
            else:
                # Regular bullet point
                p.text = f"▸  {stripped}"
                p.alignment = PP_ALIGN.LEFT
                p.font.size = Pt(15)
                p.font.bold = False
                p.font.color.rgb = self.rgb_to_pptx(self.bg.text_primary)
                p.font.name = "Calibri"
                p.space_after = Pt(8)
                p.space_before = Pt(2)
    
    # ──────────────────────────────────────────────────────────────────────────
    # COMPLETE SLIDE BUILDER — Orchestrates all 4 steps
    # ──────────────────────────────────────────────────────────────────────────
    
    def integrate_side_image(self, slide, image_path: str) -> None:
        """Place image in its designated column layout to prevent text overlap."""
        if not image_path or not os.path.exists(image_path):
            return
        layout = self.current_layout
        try:
            # Use specific image_left coordinate if defined (e.g. content_with_image_left)
            if hasattr(layout, "image_left") and layout.image_left is not None:
                img_left = self._inches(layout.image_left)
            else:
                img_left = self._inches(layout.image_right)
            
            img_top = self._inches(layout.image_top)
            img_width = self._inches(layout.image_width)
            # Maintain professional aspect ratio / layout constraints
            slide.shapes.add_picture(image_path, img_left, img_top, width=img_width, height=self._inches(layout.image_height))
        except Exception as e:
            print(f"     ⚠ Side image integration skipped: {e}")

    def build_content_slide(self, slide, title: str, content: str, slide_index: int,
                            layout_name: str = None, bg_image_path: str = None) -> None:
        """
        Build a complete slide by executing all 4 design steps in order.
        
        Step 1: Apply background design
        Step 2: Set up layout architecture
        Step 3: Integrate visual assets (image + decorative elements)
        Step 4: Format content text
        """
        # STEP 1: Background Design
        self.apply_background(slide)
        
        # STEP 2: Layout Architecture
        if layout_name and layout_name in LAYOUTS:
            self.set_layout(LAYOUTS[layout_name])
        else:
            self.set_layout(select_best_layout(slide_index, 15, has_image=(bg_image_path is not None)))
        
        self.add_header_bar(slide, title)
        self.add_accent_line(slide)
        
        # STEP 3: Visual Assets Integration
        if bg_image_path:
            layout_name = self.current_layout.name
            if layout_name in ["content_with_image_right", "content_with_image_left", "split_comparison"]:
                self.integrate_side_image(slide, bg_image_path)
            else:
                self.integrate_background_image(slide, bg_image_path)
        self.add_decorative_element(slide)
        
        # STEP 4: Content Formatting
        self.format_content_text(slide, title, content, slide_index)
    
    def build_title_slide(self, slide, title: str, subtitle: str, bg_image_path: str = None) -> None:
        """
        Build the title slide following all 4 steps.
        """
        layout = LAYOUTS["title_slide"]
        self.set_layout(layout)
        
        # STEP 1: Background Design
        self.apply_background(slide)
        
        # STEP 3: Visual Assets Integration (background image)
        if bg_image_path:
            self.integrate_background_image(slide, bg_image_path, overlay_opacity=45)
        
        # Add decorative elements
        self.add_decorative_element(slide)
        
        # Add another decorative circle on the right
        circle2 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            self._inches(self._slide_width_inches - 2.0),
            self._inches(0.5),
            self._inches(1.5), self._inches(1.5),
        )
        _apply_blended_fill(
            circle2,
            self.bg.accent_color_2[0], self.bg.accent_color_2[1], self.bg.accent_color_2[2],
            90,  # 10% opacity
        )
        _remove_shape_border(circle2)
        sp = circle2.shape._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)
        
        # Center accent bar
        center_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self._inches(4.5), self._inches(2.4),
            self._inches(4.3), Pt(4),
        )
        _apply_solid_fill(center_bar, self.bg.accent_color_1[0], self.bg.accent_color_1[1], self.bg.accent_color_1[2])
        _remove_shape_border(center_bar)
        
        # Title text
        tbox = slide.shapes.add_textbox(
            self._inches(layout.title_center_x - layout.title_width / 2),
            self._inches(layout.title_top),
            self._inches(layout.title_width),
            self._inches(1.8),
        )
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(42)
        p.font.bold = True
        p.font.color.rgb = self.rgb_to_pptx(self.bg.text_primary)
        p.font.name = "Calibri"
        
        # Subtitle text
        if subtitle:
            sbox = slide.shapes.add_textbox(
                self._inches(layout.title_center_x - layout.subtitle_width / 2),
                self._inches(layout.subtitle_top),
                self._inches(layout.subtitle_width),
                self._inches(1.2),
            )
            stf = sbox.text_frame
            stf.word_wrap = True
            sp = stf.paragraphs[0]
            sp.text = subtitle
            sp.alignment = PP_ALIGN.CENTER
            sp.font.size = Pt(20)
            sp.font.color.rgb = self.rgb_to_pptx(self.bg.text_secondary)
            sp.font.name = "Calibri"


# ══════════════════════════════════════════════════════════════════════════════
# PROFESSIONAL PRESENTATION CREATOR — Full 4-Step Orchestrator
# ══════════════════════════════════════════════════════════════════════════════

class ProfessionalPresentationCreator:
    """
    Creates presentations following the exact 4-step design process:
    
    Step 1: Background Design — Auto-select color palette based on topic
    Step 2: Professional Layout Architecture — Grid-based slide layouts
    Step 3: Visual Assets Integration — Pexels images + decorative elements
    Step 4: Content Generation & Formatting — Structured, impactful text
    """
    
    def __init__(self):
        self.pexels_client = get_pexels_client()
    
    def create_professional_presentation(
        self,
        title: str,
        slides_content: List[Dict],
        filename: str = None,
        show_progress: bool = True,
        force_background: str = None,
    ) -> str:
        """
        Create a presentation following the 4-step professional process.
        
        Args:
            title: Presentation title
            slides_content: List of slide dicts with title, content, image_query
            filename: Optional output filename
            show_progress: Print progress to console
            force_background: Force a specific background design name
        """
        # ══════════════════════════════════════════════════════════════════════
        # STEP 1: BACKGROUND DESIGN
        # ══════════════════════════════════════════════════════════════════════
        if show_progress:
            print(f"\n  {'=' * 60}")
            print(f"  STEP 1/4: BACKGROUND DESIGN")
            print(f"  Analyzing topic for optimal color palette...")
            print(f"  {'=' * 60}")
            time.sleep(0.5)
        
        if force_background and force_background in PROFESSIONAL_BACKGROUNDS:
            bg_design = PROFESSIONAL_BACKGROUNDS[force_background]
        else:
            bg_design = auto_select_background(title)
        
        if show_progress:
            print(f"  🎨 Selected: {bg_design.name}")
            print(f"  📝 {bg_design.description}")
            print(f"     Base: RGB{bg_design.base_color}")
            print(f"     Accent: RGB{bg_design.accent_color_1}")
            time.sleep(0.8)
        
        # Initialize the designer with the selected background
        designer = ProfessionalSlideDesigner(bg_design)
        
        if filename is None:
            ts = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"professional_pres_{ts}.pptx"
        if not filename.endswith(".pptx"):
            filename += ".pptx"
        filepath = os.path.join(OUTPUT_DIR, filename)
        
        # ══════════════════════════════════════════════════════════════════════
        # STEP 2: LAYOUT ARCHITECTURE + STEP 3: VISUAL ASSETS + STEP 4: CONTENT
        # ══════════════════════════════════════════════════════════════════════
        prs = Presentation()
        prs.slide_width = designer.slide_width
        prs.slide_height = designer.slide_height
        
        # ── TITLE SLIDE ──────────────────────────────────────────────────────
        if show_progress:
            print(f"\n  ─── TITLE SLIDE ───")
            print(f"  Applying layout architecture + visual assets...")
            time.sleep(0.3)
        
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Get background image for title
        bg_image = download_background_image(f"{title} background", 0)
        
        subtitle = (
            f"A Comprehensive Professional Overview\n"
            f"{datetime.now().strftime('%B %Y')}"
        )
        
        designer.build_title_slide(title_slide, title, subtitle, bg_image)
        
        if show_progress:
            print(f"  ✓ Title slide created with {bg_design.name} design")
            time.sleep(0.5)
        
        # ── CONTENT SLIDES ───────────────────────────────────────────────────
        if show_progress:
            print(f"\n  {'=' * 60}")
            print(f"  STEPS 2-4: LAYOUT + VISUALS + CONTENT")
            print(f"  Creating {len(slides_content)} content slides")
            print(f"  {'=' * 60}")
            time.sleep(0.8)
        
        for i, slide_data in enumerate(slides_content, 1):
            slide_title = slide_data.get("title", f"Slide {i}")
            slide_content = slide_data.get("content", "")
            image_query = slide_data.get("image_query", slide_title)
            
            if show_progress:
                print(f"\n  ═══ Slide {i}/{len(slides_content)} ═══")
                print(f"  📄 {slide_title}")
                time.sleep(0.3)
            
            # Get background image from Pexels
            if show_progress:
                print(f"     Step 3: Searching Pexels for '{image_query}'...")
                time.sleep(0.2)
            
            bg_image = download_background_image(image_query, i)
            
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Run all 4 steps to build the slide
            designer.build_content_slide(
                slide, slide_title, slide_content, i,
                bg_image_path=bg_image,
            )
            
            if show_progress:
                print(f"     ✓ All 4 steps complete for slide {i}")
                time.sleep(0.4)
        
        # ── SAVE ─────────────────────────────────────────────────────────────
        if show_progress:
            print(f"\n  {'=' * 60}")
            print(f"  💾 Saving professional presentation...")
            time.sleep(0.3)
        
        prs.save(filepath)
        
        if show_progress:
            print(f"  ✓ File saved: {filepath}")
            print(f"\n  ✅ PRESENTATION COMPLETE!")
            print(f"  📊 Total: {len(slides_content) + 1} slides")
            print(f"  🎨 Design: {bg_design.name}")
            print(f"  📁 File: {filename}")
            time.sleep(0.5)
        
        # Auto-open
        self._auto_open_file(filepath)
        return filepath
    
    def _auto_open_file(self, filepath: str):
        """Open the generated file with the OS default application."""
        try:
            if os.name == "nt":
                os.startfile(filepath)
            elif hasattr(os, "uname") and os.uname().sysname == "Darwin":
                import subprocess
                subprocess.Popen(["open", filepath])
            else:
                import subprocess
                subprocess.Popen(["xdg-open", filepath])
            print(f"  🚀 Presentation opened: {os.path.basename(filepath)}")
        except Exception as e:
            print(f"  [Auto-open failed]: {e}")
            print(f"  📁 Open manually: {filepath}")


# ══════════════════════════════════════════════════════════════════════════════
# PUBLIC CONVENIENCE FUNCTION
# ══════════════════════════════════════════════════════════════════════════════

def create_professional_presentation(
    title: str,
    slides_content: List[Dict],
    filename: str = None,
    show_progress: bool = True,
    force_background: str = None,
) -> str:
    """
    Create a presentation following the full 4-step professional design process.
    
    Step 1: Background Design — Auto-selects optimal color palette
    Step 2: Professional Layout Architecture — Grid-based layouts
    Step 3: Visual Assets Integration — Pexels images + decorative elements
    Step 4: Content Generation & Formatting — Structured, impactful text
    
    Args:
        title: Presentation title
        slides_content: List of slide dicts with title, content, image_query
        filename: Optional output filename
        show_progress: Show progress in console
        force_background: Force a background design by name
        
    Returns:
        Path to the saved .pptx file
    """
    creator = ProfessionalPresentationCreator()
    return creator.create_professional_presentation(
        title, slides_content, filename, show_progress, force_background
    )


# ══════════════════════════════════════════════════════════════════════════════
# PRESENTATION CONTENT GENERATOR — Step 4 helper for LLM integration
# ══════════════════════════════════════════════════════════════════════════════

# This is the prompt template that tells the LLM HOW to generate content
# following the 4-step process. It gets embedded in the system prompt.
PROFESSIONAL_PRESENTATION_SYSTEM_PROMPT = """
PRESENTATION DESIGN PROCESS — FOLLOW THESE STEPS EXACTLY IN ORDER:

STEP 1: BACKGROUND DESIGN
- First, determine the visual foundation: color palette, gradient, and texture.
- The system AUTO-SELECTS the best background based on the topic.
- Available backgrounds: midnight_professional (finance/law), slate_modern (tech/startup), 
  ivory_elegance (fashion/luxury), forest_depth (environment/sustainability),
  sunset_corporate (marketing/creative), ocean_clarity (healthcare/education).
- You do NOT choose the background — the system does it automatically from the topic.

STEP 2: PROFESSIONAL LAYOUT ARCHITECTURE
- The system uses a grid-based layout system with multiple templates:
  * content_full — Full-width text, best for detailed explanations
  * content_with_image_right — Text left, image right (60/40 split)
  * content_with_image_left — Image left, text right (40/60 split)
  * numbered_list — Large number badges, great for sequential concepts
  * split_comparison — Two-column comparison layout
  * data_visualization — Extra space for charts/data
- Layouts cycle automatically for visual variety — you don't select them.

STEP 3: VISUAL ASSETS INTEGRATION
- Each slide gets a background image from Pexels (via the image_query field).
- Decorative geometric elements are added automatically.
- Images get a semi-transparent overlay for text readability.

STEP 4: CONTENT GENERATION & FORMATTING (YOUR JOB)
Write CONCISE, IMPACTFUL, PROFESSIONAL content following these rules:
1. Minimum 10 content slides (plus title = 11 total slides)
2. Each slide title must be descriptive and specific
3. Each slide must have 4-6 bullet points of content
4. Each bullet must be a COMPLETE sentence (15+ words) with real facts
5. Use professional, authoritative language — no filler phrases
6. Include specific facts, statistics, or examples in each slide
7. Each slide must have an image_query (2-4 word Pexels search term)
8. Content must be concise — bullets should be scannable and clear

REQUIRED SLIDE STRUCTURE (follow this order):
Slide 1: Comprehensive Introduction & Background
Slide 2: Historical Context & Evolution Over Time
Slide 3: Core Concepts & Fundamental Definitions
Slide 4: Key Components & System Architecture
Slide 5: How It Works — Technical Deep Dive
Slide 6: Real-World Applications & Industry Use Cases
Slide 7: Key Benefits & Competitive Advantages
Slide 8: Challenges, Limitations & Risk Factors
Slide 9: Current Industry Trends & Market Statistics
Slide 10: Future Outlook & Emerging Developments
Slide 11: Strategic Recommendations & Final Conclusion

EXAMPLE OF GOOD CONTENT:
{
    "title": "Core Machine Learning Concepts & Algorithms",
    "content": "Supervised learning algorithms analyze labeled training data to make accurate predictions on unseen datasets with measurable performance metrics.\nUnsupervised learning techniques discover hidden patterns and natural groupings within unlabeled data through clustering and dimensionality reduction methods.\nReinforcement learning enables autonomous agents to learn optimal decision-making policies through trial-and-error interactions with dynamic environments.\nDeep neural networks with multiple hidden layers can approximate complex nonlinear functions, achieving breakthrough results in image recognition and natural language processing.",
    "image_query": "machine learning algorithm"
}

EXAMPLE OF BAD CONTENT (DO NOT GENERATE):
{
    "title": "Introduction",
    "content": "ML is useful.\nIt has benefits.\nKey concepts.",
    "image_query": "technology"
}

THE SYSTEM HANDLES:
- Background selection and gradient application (Step 1)
- Layout grid and positioning (Step 2)
- Image downloading and placement, decorative elements (Step 3)

YOUR JOB (as the LLM) IS STEP 4:
- Generate the slide titles, bullet content, and image queries
- Follow the 11-slide structure above
- Every slide must have: title, content (4-6 bullets with \\n), image_query
"""


# ══════════════════════════════════════════════════════════════════════════════
# TEST RUNNER
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    print("Professional Presentation Engine v4.0 — 4-Step Design Process")
    print("=" * 60)
    
    # Test slides (simulating what the LLM would generate in Step 4)
    test_slides = [
        {
            "title": "Introduction to Artificial Intelligence",
            "content": (
                "Artificial Intelligence represents the simulation of human intelligence "
                "processes by computer systems, enabling machines to learn from experience and data.\n"
                "The field encompasses machine learning, deep learning, neural networks, "
                "and natural language processing as its primary technical disciplines.\n"
                "AI systems today power critical applications across healthcare, finance, "
                "transportation, and manufacturing with unprecedented efficiency.\n"
                "The global AI market is projected to reach $1.8 trillion by 2030, "
                "reflecting its transformative impact across every major industry sector.\n"
                "Modern AI has evolved from rule-based expert systems to sophisticated "
                "deep learning models that can outperform humans in specific tasks."
            ),
            "image_query": "artificial intelligence future",
        },
        {
            "title": "Historical Evolution of AI Technology",
            "content": (
                "The concept of artificial intelligence dates back to the 1950s when Alan Turing "
                "proposed the famous Turing Test as a measure of machine intelligence.\n"
                "Early AI research focused on symbolic reasoning and expert systems, achieving "
                "notable successes in constrained problem domains during the 1960s and 1970s.\n"
                "The AI winter periods of the 1970s and 1980s saw reduced funding and interest "
                "due to unfulfilled promises and technological limitations of the era.\n"
                "The renaissance of AI began in the 2010s driven by three key factors: big data "
                "availability, GPU-accelerated computing, and breakthrough deep learning algorithms.\n"
                "Recent advances in transformer architectures and large language models have "
                "revolutionized natural language processing capabilities beyond previous expectations."
            ),
            "image_query": "technology history timeline",
        },
        {
            "title": "Core Concepts and Fundamental Definitions",
            "content": (
                "Machine learning algorithms are categorized into three primary paradigms: "
                "supervised, unsupervised, and reinforcement learning based on training approach.\n"
                "Neural networks consist of interconnected layers of artificial neurons that "
                "learn to recognize patterns through iterative weight adjustment during training.\n"
                "Deep learning utilizes multi-layered neural architectures capable of automatically "
                "extracting hierarchical features from raw data without manual engineering.\n"
                "Natural language processing enables computers to understand, interpret, and "
                "generate human language with applications in translation and conversation.\n"
                "Computer vision systems can analyze and interpret visual information from the "
                "world, enabling applications from facial recognition to autonomous navigation."
            ),
            "image_query": "neural network visualization",
        },
    ]
    
    # Test with auto-selected background
    print(f"\n{'─' * 60}")
    print(f"TEST: Professional Presentation Engine")
    bg = auto_select_background("Artificial Intelligence Technology")
    print(f"AUTO-SELECTED BACKGROUND: {bg.name}")
    print(f"{'─' * 60}")
    
    result = create_professional_presentation(
        "Artificial Intelligence: A Comprehensive Overview",
        test_slides,
        show_progress=True,
    )
    print(f"\n✓ Saved to: {result}")